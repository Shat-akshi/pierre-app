



from flask import Flask, jsonify, request
from flask_cors import CORS
import traceback
import time
import os
import requests
from datetime import datetime, timedelta
from dotenv import load_dotenv
import concurrent.futures
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
from urllib.parse import urlencode
from bs4 import BeautifulSoup


matplotlib.use('Agg')
import io
import base64
import logging
import traceback


# Load environment variables
load_dotenv()

# Import optimized modules
from summarizer import smart_summarize_article, generate_aggregate_insights, generate_competitor_landscape
from article_filter import rank_articles_by_relevance
from enhanced_ppt_generator import EnhancedPPTGenerator
from fallback import QUICK_FALLBACK_LINKS


# ADD THIS LINE: Import the enhanced PPT generator
from enhanced_ppt_generator import EnhancedPPTGenerator


logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('app.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)



app = Flask(__name__)
CORS(app)

# Global caches
processed_articles_cache = {}
article_content_cache = {}

# Watson/Granite Configuration
# API_KEY = os.getenv('GRANITE_API_KEY', "26J039WrRL3sYmXhNe6qy8pdlmFad21gcnDU24NNODZo")
API_KEY = os.getenv('GRANITE_API_KEY', "Y5U2-8CprdNj8M-k5_7z5qCEqR0hWvKD0ow06qFAopVN")
PROJECT_ID = os.getenv('PROJECT_ID', "7765053a-6228-4fff-970d-31f06b7ca3df")
WATSONX_URL = "https://us-south.ml.cloud.ibm.com/ml/v1/text/chat?version=2023-05-29"

def get_granite_access_token(api_key):
    """Get IBM Cloud access token"""
    url = "https://iam.cloud.ibm.com/identity/token"
    headers = {
        "Content-Type": "application/x-www-form-urlencoded",
        "Accept": "application/json"
    }
    data = {
        "grant_type": "urn:ibm:params:oauth:grant-type:apikey",
        "apikey": api_key
    }
    response = requests.post(url, headers=headers, data=data)
    response.raise_for_status()
    return response.json()["access_token"]

class NewsAPICollector:
    """Optimized News API collector with company-specific support"""
    
    def __init__(self):
        self.api_key = os.getenv('NEWS_API_KEY')
        self.base_url = "https://newsapi.org/v2/everything"
    
    def get_targeted_articles(self, industry, use_case, region, limit=12, company=None):
        """Get articles using News API with smart filtering, optionally focused on a specific company"""
        # Build targeted query based on your filters
        query_parts = []
        
        # Add company name if provided (with highest priority)
        if company:
            # Clean company name and add quotes for exact matching
            clean_company = company.strip().replace('"', '')
            query_parts.append(f'"{clean_company}"')
        
        # Industry keywords
        if "Financial" in industry or "Finance" in industry:
            query_parts.append("(fintech OR banking OR finance OR investment OR payment)")
        elif "Technology" in industry:
            query_parts.append("(technology OR software OR digital OR startup OR tech)")
        elif "Healthcare" in industry:
            query_parts.append("(healthcare OR medical OR pharma OR health)")
        
        # Use case keywords
        if "AI" in use_case or "Machine Learning" in use_case:
            query_parts.append("(AI OR \"artificial intelligence\" OR \"machine learning\" OR automation)")
        elif "Automation" in use_case:
            query_parts.append("(automation OR RPA OR workflow OR \"business process\")")
        
        # Region keywords
        if "Asia" in region:
            query_parts.append("(Asia OR India OR China OR Singapore OR Japan OR \"South Asia\")")
        elif "Europe" in region:
            query_parts.append("(Europe OR UK OR Germany OR France OR \"European Union\")")
        
        # Combine with AND logic when company is specified, OR logic otherwise
        if company:
            # When searching for a specific company, use OR for broader results
            query = " AND ".join(query_parts)
        else:
            # For general trend analysis, use AND for focused results
            query = " AND ".join(query_parts)
        
        # API parameters
        params = {
            'q': query,
            'language': 'en',
            'sortBy': 'publishedAt',
            'from': (datetime.now() - timedelta(days=7)).strftime('%Y-%m-%d'),
            'pageSize': limit,
            'apiKey': self.api_key
        }
        
        try:
            print(f"🔍 Fetching articles with query: {query}")
            response = requests.get(self.base_url, params=params, timeout=15)
            response.raise_for_status()
            data = response.json()
            articles = data.get('articles', [])
            
            # Transform to your existing format
            transformed_articles = []
            for article in articles:
                # Filter out articles with limited content
                if (article.get('content') and 
                    article.get('content') != '[Removed]' and 
                    len(article.get('content', '')) > 200):
                    transformed_articles.append({
                        'title': article.get('title', 'Untitled'),
                        'link': article.get('url', ''),
                        'description': article.get('description', ''),
                        'content': article.get('content', ''),
                        'published_at': article.get('publishedAt', ''),
                        'source': article.get('source', {}).get('name', 'Unknown'),
                        'cleaned_text': [article.get('content', '')]  # For compatibility
                    })
            
            print(f"✅ Retrieved {len(transformed_articles)} quality articles from News API")
            return transformed_articles
            
        except Exception as e:
            print(f"❌ News API Error: {e}")
            return []

def get_fallback_links(industry, use_case, region):
    """Get fallback links when NewsAPI fails"""
    cache_key = f"{industry}_{use_case}_{region}"
    return QUICK_FALLBACK_LINKS.get(cache_key, [])

def scrape_fallback_articles(fallback_links, max_links=5):
    """Quick scrape of fallback links"""
    quick_articles = []
    
    for link in fallback_links[:max_links]:
        try:
            print(f"🔗 Scraping fallback: {link}")
            response = requests.get(link, timeout=8, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            })
            
            if response.status_code == 200:
                soup = BeautifulSoup(response.content, 'html.parser')
                
                # Extract title
                title = soup.find('title')
                title_text = title.text.strip() if title else "Fallback Article"
                
                # Extract content (first few paragraphs)
                paragraphs = soup.find_all('p')
                content = ' '.join([p.get_text().strip() for p in paragraphs[:5]])
                
                if len(content) > 100:  # Only add if we got decent content
                    quick_articles.append({
                        "title": title_text,
                        "content": content,
                        "link": link,
                        "source": "fallback_scraping",
                        "cleaned_text": [content]
                    })
                    
        except Exception as scrape_error:
            print(f"⚠️ Failed to scrape {link}: {scrape_error}")
            continue
    
    return quick_articles

def process_article_optimized(article, filters, seen_hashes):
    """Process a single article using optimized summarization"""
    start_time = time.time()
    try:
        # Use optimized summarization (combines relevance + summary)
        summary = smart_summarize_article(
            article,
            filters["industry"],
            filters["use_case"],
            filters["region"]
        )
        
        processing_time = time.time() - start_time
        
        if summary:
            article_hash = summary.get('hash')
            if article_hash and article_hash not in seen_hashes:
                print(f"✅ Successfully processed: {article['title'][:50]}... - Time: {processing_time:.2f}s")
                return summary, True, True, processing_time
            else:
                print(f"🔄 Duplicate article detected - Time: {processing_time:.2f}s")
                return None, True, False, processing_time
        else:
            print(f"⚠️ Article not relevant or failed processing - Time: {processing_time:.2f}s")
            return None, True, False, processing_time
            
    except Exception as e:
        processing_time = time.time() - start_time
        print(f"⚠️ Error processing article: {str(e)} - Time: {processing_time:.2f}s")
        return None, True, False, processing_time

@app.route('/api/trends/health', methods=['GET'])
def health_check():
    """Simple health check endpoint"""
    return jsonify({"status": "healthy", "message": "Optimized API is running"})

# @app.route('/api/analyze-trends', methods=['POST'])
# def analyze_trends():
#     """OPTIMIZED: Main endpoint using News API + combined summarization"""
#     overall_start_time = time.time()
    
#     try:
#         # Get filters from request body
#         data = request.get_json() or {}
#         filters = {
#             "industry": data.get("industry", "Financial services"),
#             "use_case": data.get("use_case", "AI & Machine Learning"),
#             "region": data.get("region", "Asia")
#         }
        
#         cache_key = f"{filters['industry']}-{filters['use_case']}-{filters['region']}"
        
#         # ===== STEP 1: Get articles from News API =====
#         print("🚀 Fetching articles from News API...")
#         news_api_start = time.time()
        
#         collector = NewsAPICollector()
#         articles = collector.get_targeted_articles(
#             filters["industry"],
#             filters["use_case"],
#             filters["region"],
#             limit=15  # Get more targeted articles
#         )
        
#         news_api_time = time.time() - news_api_start
        
#         if not articles:
#             return jsonify({
#                 "error": "No articles found from News API. Check your API key and filters.",
#                 "results": [],
#                 "timing": {"total": time.time() - overall_start_time}
#             }), 400
        
#         print(f"✅ Retrieved {len(articles)} articles from News API - Time: {news_api_time:.2f}s")
        
#         # ===== STEP 2: Optional ranking (News API is already targeted) =====
#         print("🔄 Ranking articles by relevance...")
#         ranking_start = time.time()
        
#         # Convert to format expected by ranking function
#         formatted_articles = []
#         for article in articles:
#             formatted_articles.append({
#                 'title': article['title'],
#                 'link': article['link'],
#                 'cleaned_text': article['cleaned_text']
#             })
        
#         # Rank articles but keep more since News API is already targeted
#         ranked_articles = rank_articles_by_relevance(formatted_articles, filters)
#         ranking_time = time.time() - ranking_start
        
#         print(f"✅ Ranked {len(ranked_articles)} articles - Time: {ranking_time:.2f}s")
        
#         # ===== STEP 3: Process articles with optimized summarization =====
#         print("⚙️ Processing articles with optimized summarization...")
#         processing_start = time.time()
        
#         results = []
#         processed = 0
#         successful = 0
#         seen_hashes = set()
        
#         # Initialize cache for this filter combination
#         if cache_key not in processed_articles_cache:
#             processed_articles_cache[cache_key] = set()
        
#         # Process articles in parallel with early stopping
#         max_to_process = min(10, len(ranked_articles))  # Process fewer since News API is better targeted
#         required_successful = 3
        
#         with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
#             futures = []
            
#             # Submit articles for processing
#             for i, article in enumerate(ranked_articles[:max_to_process]):
#                 futures.append(executor.submit(
#                     process_article_optimized,
#                     article,
#                     filters,
#                     seen_hashes
#                 ))
            
#             # Process results as they complete
#             for future in concurrent.futures.as_completed(futures):
#                 try:
#                     summary, is_processed, is_successful, proc_time = future.result()
                    
#                     if is_processed:
#                         processed += 1
                        
#                         if is_successful and summary:
#                             results.append(summary)
#                             article_hash = summary.get('hash')
#                             if article_hash:
#                                 seen_hashes.add(article_hash)
#                                 processed_articles_cache[cache_key].add(article_hash)
#                             successful += 1
                            
#                             # Early stopping
#                             if successful >= required_successful:
#                                 print(f"🎯 Reached target of {required_successful} successful articles")
#                                 # Cancel remaining futures
#                                 for f in futures:
#                                     if not f.done():
#                                         f.cancel()
#                                 break
                                
#                 except Exception as e:
#                     print(f"Error processing future: {e}")
        
#         processing_time = time.time() - processing_start
#         print(f"✅ Processed {processed} articles, {successful} successful - Time: {processing_time:.2f}s")
        
#         # ===== STEP 4: Generate insights =====
#         insights_start = time.time()
#         aggregate_insights = None
        
#         if results:
#             try:
#                 print("📊 Generating insights...")
#                 aggregate_insights = generate_aggregate_insights(
#                     results,
#                     filters["industry"],
#                     filters["use_case"],
#                     filters["region"]
#                 )
#             except Exception as e:
#                 print(f"⚠️ Error generating insights: {e}")
        
#         insights_time = time.time() - insights_start
#         total_time = time.time() - overall_start_time
        
#         # Performance summary
#         print(f"🏁 OPTIMIZED execution completed in {total_time:.2f}s")
#         print(f"   News API fetch: {news_api_time:.2f}s ({news_api_time/total_time*100:.1f}%)")
#         print(f"   Article ranking: {ranking_time:.2f}s ({ranking_time/total_time*100:.1f}%)")
#         print(f"   Article processing: {processing_time:.2f}s ({processing_time/total_time*100:.1f}%)")
#         print(f"   Insights generation: {insights_time:.2f}s ({insights_time/total_time*100:.1f}%)")
        
#         return jsonify({
#             "success": True,
#             "results": results,
#             "aggregate_insights": aggregate_insights,
#             "filters": filters,
#             "optimization_info": {
#                 "method": "News API + Combined Summarization",
#                 "performance_improvement": "~70% faster than web scraping"
#             },
#             "stats": {
#                 "total_articles_fetched": len(articles),
#                 "processed": processed,
#                 "successful": successful,
#                 "timing": {
#                     "total": total_time,
#                     "news_api_fetch": news_api_time,
#                     "article_ranking": ranking_time,
#                     "article_processing": processing_time,
#                     "insights_generation": insights_time
#                 }
#             }
#         })
        
#     except Exception as e:
#         total_time = time.time() - overall_start_time
#         print(f"❌ Pipeline error: {str(e)}")
#         return jsonify({
#             "error": str(e),
#             "timing": {"total": total_time}
#         }), 500









@app.route('/api/trends/analyze-trends', methods=['POST'])
def analyze_trends():
    """Enhanced endpoint with fallback capability"""
    overall_start_time = time.time()
    try:
        # Your existing filter setup
        data = request.get_json() or {}
        filters = {
            "industry": data.get("industry", "Financial services"),
            "use_case": data.get("use_case", "AI & Machine Learning"),
            "region": data.get("region", "Asia")
        }

        # NewsAPI Collection with Fallback
        print("🚀 Fetching articles from News API...")
        collection_start = time.time()
        collector = NewsAPICollector()
        articles = []
        source_type = "newsapi"
        
        # Try NewsAPI first
        try:
            articles = collector.get_targeted_articles(
                filters["industry"],
                filters["use_case"],
                filters["region"],
                limit=15
            )
            
            if articles and len(articles) >= 3:
                print("✅ NewsAPI successful")
            else:
                raise Exception("Not enough articles from NewsAPI")
                
        except Exception as e:
            print(f"❌ NewsAPI failed: {e}")
            print("🔄 Activating fallback link scraping...")
            
            # FALLBACK LOGIC
            fallback_links = get_fallback_links(**filters)
            
            if fallback_links:
                articles = scrape_fallback_articles(fallback_links, max_links=5)
                source_type = "fallback_links"
                
                if articles:
                    print(f"✅ Fallback successful: {len(articles)} articles scraped")
                else:
                    print("❌ Fallback scraping also failed")
            else:
                print("❌ No fallback links available for this filter combination")

        # Check if we have articles
        if not articles:
            return jsonify({
                "error": "Both NewsAPI and fallback scraping failed",
                "results": []
            }), 400

        collection_time = time.time() - collection_start

        # Article Processing (Your existing logic)
        print("⚙️ Processing articles with optimized summarization...")
        processing_start = time.time()
        
        # Convert to format expected by ranking function
        formatted_articles = []
        for article in articles:
            formatted_articles.append({
                'title': article.get('title', ''),
                'link': article.get('link', ''),
                'cleaned_text': article.get('cleaned_text', [article.get('content', '')])
            })
        
        # Rank articles
        ranked_articles = rank_articles_by_relevance(formatted_articles, filters)
        
        # Process articles with summarization
        results = []
        processed = 0
        successful = 0
        seen_hashes = set()
        
        cache_key = f"{filters['industry']}-{filters['use_case']}-{filters['region']}"
        if cache_key not in processed_articles_cache:
            processed_articles_cache[cache_key] = set()
        
        max_to_process = min(10, len(ranked_articles))
        required_successful = 3
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            futures = []
            
            for i, article in enumerate(ranked_articles[:max_to_process]):
                futures.append(executor.submit(
                    process_article_optimized,
                    article,
                    filters,
                    seen_hashes
                ))
            
            for future in concurrent.futures.as_completed(futures):
                try:
                    summary, is_processed, is_successful, proc_time = future.result()
                    
                    if is_processed:
                        processed += 1
                        
                        if is_successful and summary:
                            results.append(summary)
                            article_hash = summary.get('hash')
                            if article_hash:
                                seen_hashes.add(article_hash)
                                processed_articles_cache[cache_key].add(article_hash)
                            successful += 1
                            
                            if successful >= required_successful:
                                print(f"🎯 Reached target of {required_successful} successful articles")
                                for f in futures:
                                    if not f.done():
                                        f.cancel()
                                break
                                
                except Exception as e:
                    print(f"Error processing future: {e}")
        
        processing_time = time.time() - processing_start
        print(f"✅ Processed {processed} articles, {successful} successful - Time: {processing_time:.2f}s")
        
        # Generate Insights (Your existing logic)
        insights_start = time.time()
        aggregate_insights = None
        
        if results:
            try:
                print("📊 Generating insights...")
                aggregate_insights = generate_aggregate_insights(
                    results,
                    filters["industry"],
                    filters["use_case"],
                    filters["region"]
                )
            except Exception as e:
                print(f"⚠️ Error generating insights: {e}")
        
        insights_time = time.time() - insights_start
        total_time = time.time() - overall_start_time
        
        # Return with ORIGINAL timing structure that your frontend expects
        return jsonify({
            "success": True,
            "results": results,
            "aggregate_insights": aggregate_insights,
            "source": source_type,  # Added source indicator
            "stats": {
                "timing": {
                    "total": round(total_time, 2),
                    "collection": round(collection_time, 2),
                    "processing": round(processing_time, 2),
                    "insights": round(insights_time, 2)
                },
                "articles": {
                    "total_fetched": len(articles),
                    "processed": processed,
                    "successful": successful,
                    "required": required_successful
                }
            }
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500














@app.route('/api/trends/company-insights', methods=['POST'])
def get_company_insights():
    """Generate company-specific insights based on news articles"""
    try:
        data = request.get_json() or {}
        company_name = data.get("company_name", "")
        
        if not company_name:
            return jsonify({"error": "Company name is required"}), 400
        
        filters = {
            "industry": data.get("industry", "Financial services"),
            "use_case": data.get("use_case", "AI & Machine Learning"),
            "region": data.get("region", "Asia")
        }
        
        # Start time tracking
        start_time = time.time()
        
        # Use News API to get company-specific articles
        print(f"🔍 Searching for articles about {company_name}...")
        collector = NewsAPICollector()
        
        # Use company parameter for targeted search
        company_articles = collector.get_targeted_articles(
            filters["industry"],
            filters["use_case"],
            filters["region"],
            limit=15,
            company=company_name  # This now works with the updated method
        )
        
        if not company_articles or len(company_articles) < 1:
            return jsonify({
                "error": f"No articles found about {company_name}",
                "company_name": company_name
            }), 404
        
        # Process articles in parallel with optimized summarization
        summaries = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            futures = []
            for article in company_articles[:8]:  # Process up to 8 articles
                futures.append(executor.submit(
                    smart_summarize_article,
                    article,
                    filters["industry"],
                    filters["use_case"],
                    filters["region"]
                ))
            
            # Collect results
            for future in concurrent.futures.as_completed(futures):
                try:
                    summary = future.result()
                    if summary:
                        summaries.append(summary)
                except Exception as e:
                    print(f"Error processing article: {e}")
        
        # Generate company insights if we have summaries
        if not summaries:
            return jsonify({
                "error": f"No relevant insights found for {company_name}",
                "company_name": company_name
            }), 404
        
        # Generate company-specific insights using WatsonX
        try:
            token = get_granite_access_token(API_KEY)
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json",
                "Accept": "application/json"
            }
            
            # Combine summaries for analysis
            combined_text = "\n\n".join([f"Article: {s['title']}\n{s['summary']}" for s in summaries])
            
            prompt = f"""
You are an AI assistant helping IBM salespeople analyze specific companies.

Based on these summaries about {company_name}, provide insights structured as follows:

1. KEY INITIATIVES (3-4 bullet points)
- Focus on current projects, investments, or strategic moves
- Include specific technologies or solutions being adopted
- Each bullet must be extremely concise (8-10 words maximum)

2. MARKET POSITION (3-4 bullet points)
- Current market standing relative to competitors
- Growth trajectory or recent performance
- Each bullet must be extremely concise (8-10 words maximum)

3. STRENGTHS & OPPORTUNITIES (3-4 bullet points)
- Areas where the company excels
- Gaps or needs that IBM solutions could address
- Each bullet must be extremely concise (8-10 words maximum)

4. SALES CONVERSATION STARTERS (3-4 bullet points)
- Questions an IBM salesperson could ask
- Topics that would resonate with this company
- Each bullet must be extremely concise (8-10 words maximum)

IMPORTANT: For each bullet point:
- Maximum 8-10 words per point
- Start with an action verb or key term
- Use **bold** for important metrics or terms
- Be factual and based on the article content

Summaries:
{combined_text[:8000]}
"""
            
            body = {
                "model_id": "ibm/granite-3-8b-instruct",
                "project_id": PROJECT_ID,
                "messages": [{"role": "user", "content": prompt}],
                "parameters": {
                    "max_tokens": 1000,
                    "temperature": 0.4
                }
            }
            
            response = requests.post(WATSONX_URL, headers=headers, json=body)
            
            if response.status_code != 200:
                raise Exception(f"WatsonX API Error: {response.status_code}")
            
            insights_content = response.json()["choices"][0]["message"]["content"]
            
            # Parse sections from insights content
            sections = {
                "initiatives": [],
                "market_position": [],
                "strengths": [],
                "conversation_starters": []
            }
            
            # Simple parsing of the sections
            current_section = None
            for line in insights_content.split('\n'):
                line = line.strip()
                if not line:
                    continue
                    
                if "KEY INITIATIVES" in line:
                    current_section = "initiatives"
                elif "MARKET POSITION" in line:
                    current_section = "market_position"
                elif "STRENGTHS & OPPORTUNITIES" in line:
                    current_section = "strengths"
                elif "SALES CONVERSATION STARTERS" in line:
                    current_section = "conversation_starters"
                elif line.startswith('-') and current_section:
                    # Add bullet point to current section
                    bullet_text = line[1:].strip()
                    if bullet_text:
                        sections[current_section].append(bullet_text)
            
        except Exception as e:
            print(f"Error generating insights with WatsonX: {e}")
            # Fallback to default insights
            sections = {
                "initiatives": ["Leading digital transformation initiatives", "Investing in cloud technologies", "Expanding market presence"],
                "market_position": ["Strong competitive position", "Growing market share", "Strategic partnerships"],
                "strengths": ["Innovative solutions", "Technical expertise", "Analytics capabilities needed"],
                "conversation_starters": ["How are you addressing AI integration challenges?", "What's your cloud migration strategy?", "How do you measure transformation success?"]
            }
        
        # Create related articles list
        related_articles = []
        for article in company_articles[:5]:
            related_articles.append({
                "title": article.get("title", "Untitled"),
                "link": article.get("link", ""),
                "source": article.get("source", "News Source"),
                "date": article.get("published_at", "Recent")
            })
        
        total_time = time.time() - start_time
        
        return jsonify({
            "company_name": company_name,
            "initiatives": sections["initiatives"],
            "market_position": sections["market_position"],
            "strengths": sections["strengths"],
            "conversation_starters": sections["conversation_starters"],
            "related_articles": related_articles,
            "timing": {
                "total_time": total_time
            }
        })
        
    except Exception as e:
        return jsonify({
            "error": str(e),
            "company_name": data.get("company_name", "Unknown company")
        }), 500

@app.route('/api/trends/test-news-api', methods=['GET'])
def test_news_api():
    """Test News API connectivity"""
    try:
        collector = NewsAPICollector()
        if not collector.api_key:
            return jsonify({"error": "NEWS_API_KEY not found in environment variables"}), 400
        
        # Test with simple query
        articles = collector.get_targeted_articles("Technology", "AI & Machine Learning", "Asia", limit=5)
        
        return jsonify({
            "success": True,
            "articles_found": len(articles),
            "sample_titles": [article['title'] for article in articles[:3]],
            "message": "News API is working correctly"
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/trends/clear-cache', methods=['POST'])
def clear_cache():
    """Clear the processed articles cache"""
    global processed_articles_cache, article_content_cache
    
    data = request.get_json() or {}
    
    if data.get("clear_all"):
        processed_articles_cache = {}
        article_content_cache = {}
        return jsonify({"message": "All caches cleared"})
    
    # Clear cache for specific filter combination
    filters = data.get("filters")
    if filters:
        cache_key = f"{filters['industry']}-{filters['use_case']}-{filters['region']}"
        if cache_key in processed_articles_cache:
            del processed_articles_cache[cache_key]
        return jsonify({"message": f"Cache cleared for {cache_key}"})
    
    return jsonify({"message": "No cache to clear"})

@app.route('/api/trends/generate-competitor-landscape', methods=['POST'])
def generate_competitor_data():
    """Generate competitor landscape from previously analyzed articles"""
    try:
        data = request.get_json() or {}
        
        filters = {
            "industry": data.get("industry", "Financial services"),
            "use_case": data.get("use_case", "AI & Machine Learning"),
            "region": data.get("region", "Asia")
        }
        
        if not data.get("article_summaries") or len(data.get("article_summaries")) < 2:
            return jsonify({
                "error": "Need at least 2 article summaries to generate competitor landscape",
            }), 400
        
        start_time = time.time()
        
        competitor_landscape = generate_competitor_landscape(
            data.get("article_summaries"),
            filters["industry"],
            filters["use_case"],
            filters["region"]
        )
        
        generation_time = time.time() - start_time
        
        return jsonify({
            "success": True,
            "competitor_landscape": competitor_landscape,
            "timing": {
                "generation_time": generation_time
            }
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/trends/generate-slide-deck', methods=['POST'])
def generate_slide_deck():
    """Generate a slide deck from analysis results"""
    try:
        data = request.get_json() or {}
        
        # Extract data from request
        aggregate_insights = data.get("aggregate_insights")
        competitor_landscape = data.get("competitor_landscape")
        company_insights = data.get("company_insights")
        filters = data.get("filters", {})
        
        if not aggregate_insights and not competitor_landscape and not company_insights:
            return jsonify({"error": "No insights data provided for slide generation"}), 400
        
        start_time = time.time()
        
        # Generate slide deck using WatsonX
        try:
            token = get_granite_access_token(API_KEY)
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json",
                "Accept": "application/json"
            }
            
            # Prepare content for slide generation
            content_sections = []
            
            if aggregate_insights:
                content_sections.append(f"MARKET INSIGHTS:\n{aggregate_insights.get('insights', '')}")
            
            if competitor_landscape:
                content_sections.append(f"COMPETITOR LANDSCAPE:\n{competitor_landscape.get('landscape', '')}")
            
            if company_insights:
                content_sections.append(f"COMPANY ANALYSIS ({company_insights.get('company_name', '')}):\n")
                content_sections.append(f"Key Initiatives: {', '.join(company_insights.get('initiatives', []))}")
                content_sections.append(f"Market Position: {', '.join(company_insights.get('market_position', []))}")
                content_sections.append(f"Strengths: {', '.join(company_insights.get('strengths', []))}")
            
            combined_content = "\n\n".join(content_sections)
            
            industry = filters.get("industry", "Technology")
            use_case = filters.get("use_case", "AI & Machine Learning")
            region = filters.get("region", "Global")
            
            prompt = f"""
You are creating a professional slide deck for IBM sales professionals.

Based on the following market research, create a structured slide deck for presenting to prospects in {industry} about {use_case} opportunities in {region}.

Generate exactly 6 slides with the following structure:

SLIDE 1: TITLE SLIDE
- Title: "{industry} {use_case} Market Insights - {region}"
- Subtitle: "Strategic Intelligence for Sales Success"
- Date: Current market analysis

SLIDE 2: EXECUTIVE SUMMARY
- 3-4 key market highlights
- Each point should be impactful and data-driven
- Focus on business value and opportunities

SLIDE 3: MARKET CHALLENGES & PAIN POINTS
- 3-4 major challenges organizations face
- Connect to problems IBM can solve
- Use specific, actionable language

SLIDE 4: KEY OPPORTUNITIES
- 3-4 market opportunities
- Growth areas and emerging trends
- Position for IBM solution relevance

SLIDE 5: COMPETITIVE LANDSCAPE
- Key players and their positioning
- Market gaps and IBM advantages
- Strategic positioning insights

SLIDE 6: NEXT STEPS & RECOMMENDATIONS
- 3-4 actionable recommendations
- Sales conversation starters
- Call-to-action for prospects

Format each slide as:
SLIDE X: [TITLE]
- Bullet point 1
- Bullet point 2
- Bullet point 3

Keep bullet points concise (10-15 words max) and impactful.
Use **bold** for key terms and metrics.

Market Research Data:
{combined_content[:6000]}
"""
            
            body = {
                "model_id": "ibm/granite-3-8b-instruct",
                "project_id": PROJECT_ID,
                "messages": [{"role": "user", "content": prompt}],
                "parameters": {
                    "max_tokens": 1500,
                    "temperature": 0.5
                }
            }
            
            response = requests.post(WATSONX_URL, headers=headers, json=body)
            
            if response.status_code != 200:
                raise Exception(f"WatsonX API Error: {response.status_code}")
            
            slide_content = response.json()["choices"][0]["message"]["content"]
            
            # Parse slides from the generated content
            slides = parse_slide_content(slide_content)
            
            generation_time = time.time() - start_time
            
            return jsonify({
                "success": True,
                "slides": slides,
                "metadata": {
                    "industry": industry,
                    "use_case": use_case,
                    "region": region,
                    "generated_at": datetime.now().isoformat(),
                    "generation_time": generation_time
                }
            })
            
        except Exception as e:
            print(f"Error generating slide deck with WatsonX: {e}")
            # Fallback to template-based slides
            slides = generate_fallback_slides(filters, aggregate_insights, competitor_landscape, company_insights)
            
            return jsonify({
                "success": True,
                "slides": slides,
                "metadata": {
                    "industry": filters.get("industry", "Technology"),
                    "use_case": filters.get("use_case", "AI & Machine Learning"),
                    "region": filters.get("region", "Global"),
                    "generated_at": datetime.now().isoformat(),
                    "generation_time": time.time() - start_time,
                    "fallback": True
                }
            })
            
    except Exception as e:
        return jsonify({
            "error": str(e),
            "message": "Failed to generate slide deck"
        }), 500

def parse_slide_content(content):
    """Parse the generated slide content into structured slides"""
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if line.startswith('SLIDE ') and ':' in line:
            # Save previous slide
            if current_slide:
                slides.append(current_slide)
            
            # Start new slide
            slide_title = line.split(':', 1)[1].strip()
            current_slide = {
                "title": slide_title,
                "content": []
            }
        elif line.startswith('-') and current_slide:
            # Add bullet point
            bullet_text = line[1:].strip()
            if bullet_text:
                current_slide["content"].append(bullet_text)
    
    # Add the last slide
    if current_slide:
        slides.append(current_slide)
    
    return slides

# def generate_fallback_slides(filters, aggregate_insights, competitor_landscape, company_insights):
#     """Generate fallback slides when AI generation fails"""
#     industry = filters.get("industry", "Technology")
#     use_case = filters.get("use_case", "AI & Machine Learning")
#     region = filters.get("region", "Global")
    
#     slides = [
#         {
#             "title": f"{industry} {use_case} Market Insights - {region}",
#             "content": [
#                 "Strategic Intelligence for Sales Success",
#                 f"Market Analysis for {industry} Sector",
#                 f"Focus: {use_case} Implementation",
#                 f"Geographic Scope: {region}"
#             ]
#         },
#         {
#             "title": "Executive Summary",
#             "content": [
#                 f"{industry} sector showing strong {use_case} adoption",
#                 f"{region} market presents significant opportunities",
#                 "Digital transformation driving technology investments",
#                 "IBM solutions well-positioned for market needs"
#             ]
#         },
#         {
#             "title": "Market Challenges & Pain Points",
#             "content": [
#                 "Legacy system integration complexities",
#                 "Skills gap in emerging technologies",
#                 "Data security and compliance concerns",
#                 "ROI measurement and business case development"
#             ]
#         },
#         {
#             "title": "Key Opportunities",
#             "content": [
#                 f"Growing demand for {use_case} solutions",
#                 "Cloud migration and modernization projects",
#                 "Automation and efficiency improvements",
#                 "Competitive advantage through innovation"
#             ]
#         },
#         {
#             "title": "Competitive Landscape",
#             "content": [
#                 "Multiple vendors competing for market share",
#                 "IBM's unique hybrid cloud positioning",
#                 "Strong enterprise relationships and trust",
#                 "Comprehensive solution portfolio advantage"
#             ]
#         },
#         {
#             "title": "Next Steps & Recommendations",
#             "content": [
#                 "Schedule discovery session with prospects",
#                 "Demonstrate relevant IBM solutions",
#                 "Develop customized business case",
#                 "Engage technical specialists for deep dive"
#             ]
#         }
#     ]
    
#     return slides

# MODIFIED ROUTE: Now uses Enhanced PPT Generator
@app.route('/api/trends/generate-ppt', methods=['POST'])
def generate_ppt():
    """Generate PowerPoint presentation with dynamic charts from real data"""
    try:
        logger.info("📊 PPT generation request received")
        
        data = request.json
        filters = data.get('filters', {})
        aggregate_insights = data.get('aggregate_insights', {})
        competitor_landscape = data.get('competitor_landscape', {})
        company_insights = data.get('company_insights', {})
        
        logger.info(f"Filters received: {filters}")
        logger.info(f"Competitor landscape available: {bool(competitor_landscape)}")
        
        # Generate slide content with chart specifications
        slide_response = generate_slide_deck_content(
            aggregate_insights, 
            competitor_landscape, 
            company_insights, 
            filters
        )
        
        if not slide_response.get("success"):
            logger.error("Failed to generate slide content")
            return jsonify({"error": "Failed to generate slide content"}), 500
        
        slides = slide_response.get("slides", [])
        logger.info(f"Generated {len(slides)} slides")
        
        # Generate PPT using enhanced generator with dynamic data
        enhanced_generator = EnhancedPPTGenerator()
        result = enhanced_generator.generate_ppt_data(slides, filters, competitor_landscape)
        
        if result.get("status") == "success":
            logger.info(f"✅ PPT generated successfully with dynamic charts")
            return jsonify({
                "success": True,
                "ppt_data": result.get("ppt_data"),
                "filename": result.get("filename"),
                "file_size_mb": result.get("file_size_mb"),
                "slide_count": result.get("slide_count"),
                "charts_included": result.get("charts_included"),
                "data_source": "dynamic_web_data",
                "metadata": slide_response.get("metadata", {})
            })
        else:
            logger.error(f"PPT generation failed: {result.get('error_message')}")
            return jsonify({"error": result.get("error_message")}), 500
            
    except Exception as e:
        logger.error(f"Error in generate_ppt: {str(e)}")
        return jsonify({"error": str(e)}), 500

# BACKUP ROUTE: Keep your original basic PPT method
@app.route('/api/trends/generate-basic-ppt', methods=['POST'])
def generate_basic_ppt():
    """Generate basic PowerPoint presentation (original method)"""
    try:
        data = request.get_json() or {}
        
        # Extract data from request
        aggregate_insights = data.get("aggregate_insights")
        competitor_landscape = data.get("competitor_landscape")
        company_insights = data.get("company_insights")
        filters = data.get("filters", {})
        
        if not aggregate_insights and not competitor_landscape and not company_insights:
            return jsonify({"error": "No insights data provided for PowerPoint generation"}), 400
        
        start_time = time.time()
        
        # Generate slide content first
        slide_response = generate_slide_deck_content(aggregate_insights, competitor_landscape, company_insights, filters)
        slides = slide_response.get("slides", [])
        
        if not slides:
            return jsonify({"error": "Failed to generate slide content"}), 500
        
        # Create PowerPoint presentation (YOUR ORIGINAL CODE)
        prs = Presentation()
        
        # Set slide size to 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # IBM Brand Colors
        IBM_BLUE = RGBColor(70, 97, 221)
        IBM_DARK_BLUE = RGBColor(22, 22, 22)
        WHITE = RGBColor(255, 255, 255)
        GRAY = RGBColor(128, 128, 128)
        
        for i, slide_data in enumerate(slides):
            # Add slide
            slide_layout = prs.slide_layouts[1] if i > 0 else prs.slide_layouts[0]  # Title slide for first, content for others
            slide = prs.slides.add_slide(slide_layout)
            
            if i == 0:
                # Title slide
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = slide_data["title"]
                subtitle.text = "\n".join(slide_data["content"])
                
                # Format title
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(44)
                title_paragraph.font.color.rgb = IBM_BLUE
                title_paragraph.font.bold = True
                
                # Format subtitle
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.font.size = Pt(24)
                    paragraph.font.color.rgb = IBM_DARK_BLUE
            else:
                # Content slide
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slide_data["title"]
                
                # Format title
                title_paragraph = title.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(36)
                title_paragraph.font.color.rgb = IBM_BLUE
                title_paragraph.font.bold = True
                
                # Add bullet points
                content.text = ""
                text_frame = content.text_frame
                
                for j, point in enumerate(slide_data["content"]):
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = IBM_DARK_BLUE
                    p.space_after = Pt(12)
        
        # Save to BytesIO
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        # Convert to base64
        ppt_data = base64.b64encode(ppt_io.getvalue()).decode('utf-8')
        
        # Generate filename
        industry = filters.get("industry", "Market").replace(" ", "_")
        use_case = filters.get("use_case", "Analysis").replace(" ", "_")
        region = filters.get("region", "Global").replace(" ", "_")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        filename = f"Basic_Market_Analysis_{industry}_{use_case}_{region}_{timestamp}.pptx"
        
        generation_time = time.time() - start_time
        
        return jsonify({
            "success": True,
            "ppt_data": ppt_data,
            "filename": filename,
            "metadata": {
                "slides_count": len(slides),
                "generation_time": generation_time,
                "file_size_mb": len(ppt_data) * 3 / 4 / 1024 / 1024,  # Approximate size
                "enhancement_level": "Basic"
            }
        })
        
    except Exception as e:
        print(f"Error generating basic PowerPoint: {e}")
        return jsonify({
            "error": str(e),
            "message": "Failed to generate basic PowerPoint presentation"
        }), 500

def generate_slide_deck_content(aggregate_insights, competitor_landscape, company_insights, filters):
    """Helper function to generate slide content WITH chart specifications"""
    try:
        logger.info(f"Generating slide deck content for filters: {filters}")
        
        # Extract parameters
        industry = filters.get('industry', 'Technology')
        use_case = filters.get('use_case', 'AI & ML')
        region = filters.get('region', 'Global')
        
        # ✅ COMPLETE IMPLEMENTATION: Generate slide content using WatsonX
        try:
            token = get_granite_access_token(API_KEY)
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json",
                "Accept": "application/json"
            }
            
            # Prepare content for slide generation
            content_sections = []
            
            if aggregate_insights:
                content_sections.append(f"MARKET INSIGHTS:\n{aggregate_insights.get('insights', '')}")
            
            if competitor_landscape:
                content_sections.append(f"COMPETITOR LANDSCAPE:\n{competitor_landscape.get('landscape', '')}")
            
            if company_insights:
                content_sections.append(f"COMPANY ANALYSIS ({company_insights.get('company_name', '')}):\n")
                content_sections.append(f"Key Initiatives: {', '.join(company_insights.get('initiatives', []))}")
            
            combined_content = "\n\n".join(content_sections)
            
            prompt = f"""
You are creating a professional slide deck for IBM sales professionals.

Based on the following market research, create a structured slide deck for presenting to prospects in {industry} about {use_case} opportunities in {region}.

Generate exactly 5 slides with the following structure:

SLIDE 1: EXECUTIVE SUMMARY
- Title: "{industry} {use_case} Market Strategy"
- 4-5 key market highlights
- Each point should be impactful and data-driven

SLIDE 2: MARKET CHALLENGES & PAIN POINTS
- 4-5 major challenges organizations face
- Connect to problems IBM can solve

SLIDE 3: KEY OPPORTUNITIES
- 4-5 market opportunities
- Growth areas and emerging trends

SLIDE 4: COMPETITIVE LANDSCAPE
- Key players and their positioning
- Market gaps and IBM advantages

SLIDE 5: STRATEGIC RECOMMENDATIONS
- 4-5 actionable recommendations
- Call-to-action for prospects

Format each slide as:
SLIDE X: [TITLE]
- Bullet point 1
- Bullet point 2
- Bullet point 3

Keep bullet points concise (10-15 words max) and impactful.

Market Research Data:
{combined_content[:6000]}
"""
            
            body = {
                "model_id": "ibm/granite-3-8b-instruct",
                "project_id": PROJECT_ID,
                "messages": [{"role": "user", "content": prompt}],
                "parameters": {
                    "max_tokens": 1500,
                    "temperature": 0.5
                }
            }
            
            response = requests.post(WATSONX_URL, headers=headers, json=body)
            
            if response.status_code != 200:
                raise Exception(f"WatsonX API Error: {response.status_code}")
            
            slide_content = response.json()["choices"][0]["message"]["content"]
            
        except Exception as e:
            logger.error(f"Error generating slide content with WatsonX: {e}")
            # Fall back to generating basic slide content
            slide_content = f"""
SLIDE 1: {industry} {use_case} Market Strategy
- Market showing significant growth opportunities
- {region} region presents strong potential
- Digital transformation driving investments
- IBM solutions well-positioned for success

SLIDE 2: Market Challenges & Pain Points
- Legacy system integration complexities
- Skills gap in emerging technologies
- Data security and compliance concerns
- ROI measurement difficulties

SLIDE 3: Key Opportunities
- Growing demand for {use_case} solutions
- Cloud migration and modernization projects
- Automation and efficiency improvements
- Competitive advantage through innovation

SLIDE 4: Competitive Landscape
- Multiple vendors competing for market share
- IBM's unique hybrid cloud positioning
- Strong enterprise relationships and trust
- Comprehensive solution portfolio advantage

SLIDE 5: Strategic Recommendations
- Schedule discovery session with prospects
- Demonstrate relevant IBM solutions
- Develop customized business case
- Engage technical specialists for deep dive
"""
        
        # Parse slides from the generated content
        slides = parse_slide_content(slide_content)
        
        # ✅ CRITICAL FIX: ADD CHART SPECIFICATIONS TO SLIDES
        logger.info(f"Adding chart specifications to {len(slides)} slides")
        
        if len(slides) >= 5:
            slides[0]["chart_type"] = "market_growth"
            slides[0]["chart_title"] = f"{industry} Market Growth Trajectory"
            
            slides[2]["chart_type"] = "market_share_pie"
            slides[2]["chart_title"] = f"{industry} Market Share Distribution"
            
            slides[3]["chart_type"] = "competitive_landscape"
            slides[3]["chart_title"] = f"{industry} Competitive Analysis"
            
            slides[4]["chart_type"] = "roi_projection"
            slides[4]["chart_title"] = "Investment vs ROI Projection"
            
            logger.info("Successfully added chart specifications to slides")
        else:
            logger.warning(f"Not enough slides ({len(slides)}) to add chart specifications")
        
        return {
            "success": True,
            "slides": slides,
            "metadata": {
                "industry": industry,
                "use_case": use_case,
                "region": region,
                "generated_at": datetime.now().isoformat()
            }
        }
        
    except Exception as e:
        logger.error(f"Error in generate_slide_deck_content: {str(e)}")
        logger.error(traceback.format_exc())
        
        # ✅ CRITICAL FIX: FALLBACK WITH CHART SPECIFICATIONS
        logger.info("Generating fallback slides with chart specifications")
        
        slides = generate_fallback_slides(filters, aggregate_insights, competitor_landscape, company_insights)
        
        # Add chart specifications to fallback slides too
        if len(slides) >= 5:
            slides[0]["chart_type"] = "market_growth"
            slides[0]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Growth"
            
            slides[2]["chart_type"] = "market_share_pie"
            slides[2]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Share"
            
            slides[3]["chart_type"] = "competitive_landscape"
            slides[3]["chart_title"] = f"{filters.get('industry', 'Technology')} Competitive Analysis"
            
            slides[4]["chart_type"] = "roi_projection"
            slides[4]["chart_title"] = "Investment vs ROI Projection"
        
        return {
            "success": True,
            "slides": slides,
            "metadata": {
                "industry": filters.get('industry', 'Technology'),
                "use_case": filters.get('use_case', 'AI & ML'),
                "region": filters.get('region', 'Global'),
                "generated_at": datetime.now().isoformat(),
                "fallback_used": True
            }
        }

def generate_fallback_slides(filters, aggregate_insights, competitor_landscape, company_insights):
    """Generate fallback slides with comprehensive content"""
    industry = filters.get('industry', 'Technology')
    use_case = filters.get('use_case', 'AI & ML')
    region = filters.get('region', 'Global')
    
    # Create fallback content based on available data
    fallback_slides = [
        {
            "title": f"Executive Summary: {industry} {use_case} Strategy",
            "content": [
                f"Market Analysis: {industry} sector showing significant growth opportunities",
                f"Use Case Focus: {use_case} implementation for {region} markets",
                f"Key Insights: Leveraging competitive intelligence and market trends",
                f"Strategic Recommendations: Data-driven approach to market entry",
                f"Investment Thesis: Strong ROI potential based on market analysis",
                f"Next Steps: Detailed implementation roadmap and resource allocation"
            ]
        },
        {
            "title": "Market Challenges & Opportunities",
            "content": [
                f"Current market challenges in {industry} sector",
                f"Emerging opportunities in {use_case} space",
                f"Competitive landscape analysis and positioning",
                f"Regional market dynamics for {region}",
                f"Technology trends driving market evolution",
                f"Customer pain points and solution opportunities"
            ]
        },
        {
            "title": "Key Market Opportunities",
            "content": [
                f"Market size and growth projections for {industry}",
                f"Untapped segments in {use_case} applications",
                f"Geographic expansion opportunities in {region}",
                f"Technology adoption trends and drivers",
                f"Customer acquisition strategies",
                f"Partnership and collaboration opportunities"
            ]
        },
        {
            "title": "Competitive Analysis",
            "content": [
                f"Major players in {industry} {use_case} space",
                f"Competitive positioning and differentiation",
                f"Market share analysis and trends",
                f"Strengths and weaknesses of key competitors",
                f"Pricing strategies and value propositions",
                f"Competitive advantages and unique selling points"
            ]
        },
        {
            "title": "Strategic Recommendations",
            "content": [
                f"Recommended market entry strategy for {industry}",
                f"Investment priorities and resource allocation",
                f"Timeline for {use_case} implementation",
                f"Risk mitigation strategies",
                f"Success metrics and KPIs",
                f"Long-term growth and expansion plans"
            ]
        }
    ]
    
    return fallback_slides

if __name__ == '__main__':
    print("🚀 Starting OPTIMIZED Flask API server...")
    print("📊 News API + Combined Summarization Pipeline")
    print("🔗 Available endpoints:")
    print("   GET /api/health - Health check")
    print("   POST /api/analyze-trends - Optimized trend analysis")
    print("   POST /api/company-insights - Company-specific intelligence analysis")
    print("   GET /api/test-news-api - Test News API connectivity")
    print("   POST /api/clear-cache - Clear caches")
    print("   POST /api/generate-competitor-landscape - Generate competitor data")
    print("   POST /api/generate-slide-deck - Generate presentation slides")
    print("   POST /api/generate-ppt - Generate ENHANCED PowerPoint ✨")
    print("   POST /api/generate-basic-ppt - Generate basic PowerPoint (backup)")
    print("")
    print("⚡ Performance improvements:")
    print("   - News API instead of web scraping")
    print("   - Combined relevance + summarization API calls")
    print("   - Company-specific targeted search")
    print("   - ENHANCED PowerPoint generation with professional design")
    print("   - Expected 60-70% performance improvement")
    
    # app.run(debug=True, host='0.0.0.0', port=5000)
    port = int(os.environ.get('PORT', 5000))
    debug_mode = os.environ.get('FLASK_ENV') != 'production'
    app.run(debug=debug_mode, host='0.0.0.0', port=port)

