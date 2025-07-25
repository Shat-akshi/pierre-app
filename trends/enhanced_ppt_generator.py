# import io
# import base64
# import logging
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.dml import MSO_THEME_COLOR
# import matplotlib.pyplot as plt
# import matplotlib.patches as patches
# import numpy as np
# import seaborn as sns
# from datetime import datetime, timedelta
# import random
# import json
# import re

# # Configure matplotlib for better chart quality
# plt.rcParams['figure.facecolor'] = 'white'
# plt.rcParams['axes.facecolor'] = 'white'
# plt.rcParams['font.size'] = 12
# plt.rcParams['font.family'] = 'sans-serif'

# # Set seaborn style
# sns.set_style("whitegrid")
# sns.set_palette("husl")

# logger = logging.getLogger(__name__)

# class EnhancedPPTGenerator:
#     def __init__(self):
#         self.colors = {
#             'primary': RGBColor(37, 99, 235),      # Blue
#             'secondary': RGBColor(124, 58, 237),    # Purple  
#             'accent': RGBColor(5, 150, 105),       # Green
#             'warning': RGBColor(217, 119, 6),      # Orange
#             'danger': RGBColor(220, 38, 38),       # Red
#             'dark': RGBColor(31, 41, 55),          # Dark Gray
#             'light': RGBColor(249, 250, 251),      # Light Gray
#             'text': RGBColor(55, 65, 81)           # Text Gray
#         }
        
#         # Chart color palette
#         self.chart_colors = ['#2563EB', '#7C3AED', '#059669', '#D97706', '#DC2626', '#0891B2', '#C2410C']
        
#         # Industry-specific comprehensive data
#         self.industry_data = {
#             'Technology': {
#                 'market_size': 4500,
#                 'growth_rate': 12.5,
#                 'key_players': ['Microsoft', 'Google', 'Amazon', 'Apple', 'Meta'],
#                 'market_shares': [22, 18, 15, 12, 8],
#                 'growth_rates': [8, 15, 12, 6, 25],
#                 'revenue_data': [100, 120, 145, 175, 210, 250, 300, 360],
#                 'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
#                 'roi_timeline': [2024, 2025, 2026, 2027],
#                 'roi_values': [-10, 25, 85, 180],
#                 'investment_values': [5, 8, 12, 15]
#             },
#             'Healthcare': {
#                 'market_size': 3200,
#                 'growth_rate': 8.7,
#                 'key_players': ['Johnson & Johnson', 'Pfizer', 'Roche', 'Novartis', 'Merck'],
#                 'market_shares': [15, 12, 10, 9, 8],
#                 'growth_rates': [6, 8, 12, 15, 10],
#                 'revenue_data': [80, 95, 110, 125, 140, 160, 180, 205],
#                 'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
#                 'roi_timeline': [2024, 2025, 2026, 2027],
#                 'roi_values': [-15, 20, 65, 150],
#                 'investment_values': [8, 12, 18, 25]
#             },
#             'Financial Services': {
#                 'market_size': 2800,
#                 'growth_rate': 6.3,
#                 'key_players': ['JPMorgan Chase', 'Bank of America', 'Wells Fargo', 'Goldman Sachs', 'Morgan Stanley'],
#                 'market_shares': [18, 15, 12, 10, 8],
#                 'growth_rates': [4, 6, 8, 12, 15],
#                 'revenue_data': [90, 100, 110, 120, 130, 145, 160, 175],
#                 'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
#                 'roi_timeline': [2024, 2025, 2026, 2027],
#                 'roi_values': [-8, 15, 45, 120],
#                 'investment_values': [6, 10, 15, 20]
#             }
#         }

#     def normalize_industry_name(self, industry_name):
#         """Normalize industry names to match data keys"""
#         mapping = {
#             'financial services': 'Financial Services',
#             'Financial services': 'Financial Services',
#             'financial': 'Financial Services',
#             'technology': 'Technology',
#             'tech': 'Technology',
#             'healthcare': 'Healthcare',
#             'health': 'Healthcare'
#         }
#         return mapping.get(industry_name.lower(), 'Technology')

#     def extract_data_from_content(self, slides_data, filters):
#         """Extract chart data from actual slide content"""
#         extracted_data = {
#             'companies': [],
#             'market_values': [],
#             'growth_rates': [],
#             'years': [],
#             'revenue_data': []
#         }
        
#         # Parse content for numbers, companies, percentages
#         for slide in slides_data:
#             content_text = ' '.join(slide.get('content', []))
            
#             # Extract company names (look for patterns like "Company Name leads with X%")
#             company_patterns = re.findall(r'([A-Z][a-z]+ [A-Z][a-z]+|[A-Z][a-z]+) leads with (\d+)%', content_text)
#             for company, percentage in company_patterns:
#                 if company not in extracted_data['companies']:
#                     extracted_data['companies'].append(company)
#                     extracted_data['market_values'].append(int(percentage))
            
#             # Extract growth rates, years, revenue figures
#             growth_matches = re.findall(r'(\d+(?:\.\d+)?)% growth', content_text)
#             revenue_matches = re.findall(r'\$(\d+(?:\.\d+)?)B', content_text)
#             year_matches = re.findall(r'(20\d{2})', content_text)
            
#             extracted_data['growth_rates'].extend([float(x) for x in growth_matches])
#             extracted_data['revenue_data'].extend([float(x) for x in revenue_matches])
#             extracted_data['years'].extend([int(x) for x in year_matches])
        
#         return extracted_data

#     def create_market_growth_chart(self, industry):
#         """Create a professional market growth chart"""
#         try:
#             data = self.industry_data.get(industry, self.industry_data['Technology'])
            
#             fig, ax = plt.subplots(figsize=(12, 8))
#             fig.patch.set_facecolor('white')
            
#             years = data['years']
#             revenue = data['revenue_data']
            
#             # Main growth line
#             ax.plot(years, revenue, linewidth=4, color=self.chart_colors[0], 
#                    marker='o', markersize=8, label='Market Growth')
            
#             # Fill area under curve
#             ax.fill_between(years, revenue, alpha=0.2, color=self.chart_colors[0])
            
#             # Add trend line
#             z = np.polyfit(range(len(years)), revenue, 1)
#             p = np.poly1d(z)
#             ax.plot(years, p(range(len(years))), "--", color=self.chart_colors[3], 
#                    linewidth=3, alpha=0.8, label='Trend Line')
            
#             # Highlight projection period
#             ax.axvline(x=2023.5, color='gray', linestyle=':', alpha=0.7, linewidth=2)
#             ax.text(2021.5, max(revenue)*0.9, 'Historical', fontsize=12, 
#                    bbox=dict(boxstyle="round,pad=0.3", facecolor='lightblue', alpha=0.7))
#             ax.text(2025, max(revenue)*0.9, 'Projected', fontsize=12,
#                    bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgreen', alpha=0.7))
            
#             # Add data labels
#             for i, (year, val) in enumerate(zip(years, revenue)):
#                 if i % 2 == 0:  # Show every other label to avoid crowding
#                     ax.annotate(f'${val}B', (year, val), textcoords="offset points", 
#                               xytext=(0,10), ha='center', fontweight='bold')
            
#             # Styling
#             ax.set_title(f'{industry} Market Growth Trajectory (2020-2027)', 
#                         fontsize=16, fontweight='bold', pad=20)
#             ax.set_xlabel('Year', fontsize=14, fontweight='bold')
#             ax.set_ylabel('Market Value (Billions USD)', fontsize=14, fontweight='bold')
#             ax.grid(True, alpha=0.3)
#             ax.legend(fontsize=12)
            
#             # Remove top and right spines
#             ax.spines['top'].set_visible(False)
#             ax.spines['right'].set_visible(False)
            
#             plt.tight_layout()
            
#             # Save to bytes
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
#                        facecolor='white', edgecolor='none')
#             img_buffer.seek(0)
#             plt.close()
            
#             logger.info(f"Successfully created market growth chart for {industry}")
#             return img_buffer
            
#         except Exception as e:
#             logger.error(f"Error creating market growth chart: {str(e)}")
#             plt.close()
#             return None

#     def create_competitive_landscape_chart(self, industry):
#         """Create a competitive landscape chart"""
#         try:
#             data = self.industry_data.get(industry, self.industry_data['Technology'])
            
#             fig, ax = plt.subplots(figsize=(12, 8))
#             fig.patch.set_facecolor('white')
            
#             companies = data['key_players']
#             market_shares = data['market_shares']
#             growth_rates = data['growth_rates']
            
#             x = np.arange(len(companies))
#             width = 0.35
            
#             # Create bars
#             bars1 = ax.bar(x - width/2, market_shares, width, label='Market Share (%)', 
#                           color=self.chart_colors[0], alpha=0.8)
#             bars2 = ax.bar(x + width/2, growth_rates, width, label='Growth Rate (%)', 
#                           color=self.chart_colors[1], alpha=0.8)
            
#             # Add value labels on bars
#             for bar in bars1:
#                 height = bar.get_height()
#                 ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
#                        f'{height}%', ha='center', va='bottom', fontweight='bold')
            
#             for bar in bars2:
#                 height = bar.get_height()
#                 ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
#                        f'{height}%', ha='center', va='bottom', fontweight='bold')
            
#             # Styling
#             ax.set_title(f'{industry} Competitive Analysis: Market Share vs Growth Rate', 
#                         fontsize=16, fontweight='bold', pad=20)
#             ax.set_xlabel('Companies', fontsize=14, fontweight='bold')
#             ax.set_ylabel('Percentage', fontsize=14, fontweight='bold')
#             ax.set_xticks(x)
#             ax.set_xticklabels(companies, rotation=45, ha='right')
#             ax.legend(fontsize=12)
#             ax.grid(True, alpha=0.3)
            
#             # Remove top and right spines
#             ax.spines['top'].set_visible(False)
#             ax.spines['right'].set_visible(False)
            
#             plt.tight_layout()
            
#             # Save to bytes
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
#                        facecolor='white', edgecolor='none')
#             img_buffer.seek(0)
#             plt.close()
            
#             logger.info(f"Successfully created competitive landscape chart for {industry}")
#             return img_buffer
            
#         except Exception as e:
#             logger.error(f"Error creating competitive landscape chart: {str(e)}")
#             plt.close()
#             return None

#     def create_roi_projection_chart(self, industry):
#         """Create ROI projection chart"""
#         try:
#             data = self.industry_data.get(industry, self.industry_data['Technology'])
            
#             fig, ax1 = plt.subplots(figsize=(12, 8))
#             fig.patch.set_facecolor('white')
            
#             timeline = data['roi_timeline']
#             roi_values = data['roi_values']
#             investment_values = data['investment_values']
            
#             # ROI line chart
#             color1 = self.chart_colors[2]
#             ax1.set_xlabel('Year', fontsize=14, fontweight='bold')
#             ax1.set_ylabel('ROI (%)', color=color1, fontsize=14, fontweight='bold')
#             line1 = ax1.plot(timeline, roi_values, 'o-', linewidth=4, color=color1, 
#                            markersize=10, label='ROI %')
#             ax1.tick_params(axis='y', labelcolor=color1)
#             ax1.axhline(y=0, color='gray', linestyle='-', alpha=0.3)
            
#             # Investment bar chart
#             ax2 = ax1.twinx()
#             color2 = self.chart_colors[0]
#             ax2.set_ylabel('Investment ($M)', color=color2, fontsize=14, fontweight='bold')
#             bars = ax2.bar(timeline, investment_values, alpha=0.6, color=color2, 
#                           width=0.6, label='Investment ($M)')
#             ax2.tick_params(axis='y', labelcolor=color2)
            
#             # Add value labels
#             for year, roi in zip(timeline, roi_values):
#                 ax1.annotate(f'{roi}%', (year, roi), textcoords="offset points", 
#                            xytext=(0,15), ha='center', fontweight='bold', color=color1)
            
#             for year, inv in zip(timeline, investment_values):
#                 ax2.annotate(f'${inv}M', (year, inv), textcoords="offset points", 
#                            xytext=(0,5), ha='center', fontweight='bold', color=color2)
            
#             # Styling
#             ax1.set_title('Investment vs ROI Projection (2024-2027)', 
#                          fontsize=16, fontweight='bold', pad=20)
#             ax1.grid(True, alpha=0.3)
            
#             plt.tight_layout()
            
#             # Save to bytes
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
#                        facecolor='white', edgecolor='none')
#             img_buffer.seek(0)
#             plt.close()
            
#             logger.info(f"Successfully created ROI projection chart for {industry}")
#             return img_buffer
            
#         except Exception as e:
#             logger.error(f"Error creating ROI projection chart: {str(e)}")
#             plt.close()
#             return None

#     def create_market_share_pie_chart(self, industry):
#         """Create market share pie chart"""
#         try:
#             data = self.industry_data.get(industry, self.industry_data['Technology'])
            
#             fig, ax = plt.subplots(figsize=(10, 10))
#             fig.patch.set_facecolor('white')
            
#             companies = data['key_players']
#             market_shares = data['market_shares']
            
#             # Add "Others" category
#             others_share = 100 - sum(market_shares)
#             companies_with_others = companies + ['Others']
#             shares_with_others = market_shares + [others_share]
            
#             # Create pie chart
#             wedges, texts, autotexts = ax.pie(shares_with_others, labels=companies_with_others, 
#                                             autopct='%1.1f%%', startangle=90,
#                                             colors=self.chart_colors[:len(companies_with_others)],
#                                             explode=[0.05 if i == 0 else 0 for i in range(len(companies_with_others))])
            
#             # Enhance text
#             for autotext in autotexts:
#                 autotext.set_color('white')
#                 autotext.set_fontweight('bold')
#                 autotext.set_fontsize(12)
            
#             for text in texts:
#                 text.set_fontsize(11)
#                 text.set_fontweight('bold')
            
#             ax.set_title(f'{industry} Market Share Distribution', 
#                         fontsize=16, fontweight='bold', pad=20)
            
#             plt.tight_layout()
            
#             # Save to bytes
#             img_buffer = io.BytesIO()
#             plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
#                        facecolor='white', edgecolor='none')
#             img_buffer.seek(0)
#             plt.close()
            
#             logger.info(f"Successfully created market share pie chart for {industry}")
#             return img_buffer
            
#         except Exception as e:
#             logger.error(f"Error creating market share pie chart: {str(e)}")
#             plt.close()
#             return None

#     def generate_comprehensive_content(self, filters):
#         """Generate comprehensive slide content with chart specifications"""
#         industry = filters.get('industry', 'Technology')
#         use_case = filters.get('use_case', 'AI & ML')
#         region = filters.get('region', 'Global')
        
#         # Normalize industry name
#         industry = self.normalize_industry_name(industry)
#         data = self.industry_data.get(industry, self.industry_data['Technology'])
        
#         slides = [
#             {
#                 "title": f"Executive Summary: {industry} {use_case} Market Strategy",
#                 "content": [
#                     f"Market Opportunity: ${data['market_size']}B global {industry.lower()} market growing at {data['growth_rate']}% CAGR",
#                     f"Strategic Focus: {use_case} implementation targeting {region} markets with enterprise-first approach",
#                     f"Competitive Landscape: Top 5 players control 75% market share, creating consolidation opportunities",
#                     f"Investment Thesis: $15M investment over 24 months targeting 180% ROI by year 3",
#                     f"Key Value Propositions: 40% faster implementation, 60% cost reduction, 95% customer satisfaction",
#                     f"Market Entry Strategy: Phased rollout starting with pilot customers in Q2 2024",
#                     f"Revenue Projections: $2M Year 1, $12M Year 2, $35M Year 3 with 85% gross margins",
#                     f"Success Metrics: 50 enterprise clients, 8% market share, $50M ARR by 2027"
#                 ],
#                 "chart_type": "market_growth",
#                 "chart_title": f"{industry} Market Growth Trajectory"
#             },
#             {
#                 "title": f"Market Challenges & Pain Points",
#                 "content": [
#                     f"Rapid technological obsolescence requiring constant innovation cycles",
#                     f"Intense competition from both established players and disruptive startups",
#                     f"Regulatory scrutiny around data privacy and antitrust concerns",
#                     f"Talent acquisition and retention in highly competitive market",
#                     f"Cybersecurity threats and infrastructure vulnerabilities",
#                     f"Market fragmentation with 200+ vendors creating customer confusion",
#                     f"Average implementation timeline of 18-24 months causing project delays",
#                     f"Limited ROI visibility leading to 35% project abandonment rate"
#                 ],
#                 "chart_type": None,
#                 "chart_title": None
#             },
#             {
#                 "title": f"Key Market Opportunities",
#                 "content": [
#                     f"AI and machine learning integration across all business functions",
#                     f"Edge computing and 5G network expansion opportunities",
#                     f"Digital transformation acceleration post-pandemic",
#                     f"Emerging markets adoption of cloud technologies",
#                     f"Sustainability tech and green computing initiatives",
#                     f"Underserved mid-market segment representing $800M opportunity",
#                     f"Geographic expansion potential in emerging markets",
#                     f"Cross-selling opportunities with existing customer base"
#                 ],
#                 "chart_type": "market_share_pie",
#                 "chart_title": f"{industry} Market Share Distribution"
#             },
#             {
#                 "title": "Competitive Landscape Analysis",
#                 "content": [
#                     f"Market Leaders: {data['key_players'][0]} leads with {data['market_shares'][0]}% share, followed by {data['key_players'][1]} ({data['market_shares'][1]}%)",
#                     f"Growth Leaders: {data['key_players'][4]} showing highest growth at {data['growth_rates'][4]}% vs industry average {data['growth_rate']}%",
#                     f"Competitive Positioning: We target underserved mid-market with 35% better price-performance ratio",
#                     f"Differentiation Strategy: AI-first approach vs competitors' legacy modernization focus",
#                     f"Technology Advantage: Proprietary algorithms delivering 95% accuracy vs 78% industry standard",
#                     f"Customer Acquisition: Average sales cycle 6 months vs industry standard 12-18 months",
#                     f"Pricing Strategy: Premium positioning 15% above market justified by superior outcomes",
#                     f"Partnership Approach: Strategic alliances with {data['key_players'][2]} and {data['key_players'][3]} for market access"
#                 ],
#                 "chart_type": "competitive_landscape",
#                 "chart_title": f"{industry} Competitive Analysis"
#             },
#             {
#                 "title": "Strategic Recommendations & Next Steps",
#                 "content": [
#                     f"Immediate Actions (30 days): Secure Series A funding, finalize product roadmap, hire VP Engineering",
#                     f"Short-term Goals (90 days): Complete MVP development, sign 3 pilot customers, establish advisory board",
#                     f"Market Entry (6 months): Launch in {region}, build sales team of 15 FTEs, initiate marketing campaigns",
#                     f"Scale Phase (12 months): Expand to 2 additional regions, achieve $5M ARR, establish channel partnerships",
#                     f"Investment Recommendation: Approve $15M funding based on 180% projected ROI and market timing",
#                     f"Resource Allocation: 60% product development, 25% sales & marketing, 15% operations",
#                     f"Risk Mitigation: Diversify customer base, build IP portfolio, establish regulatory compliance framework",
#                     f"Success Metrics: Track monthly ARR growth, customer acquisition cost, net promoter score",
#                     f"Partnership Strategy: Engage tier-1 system integrators, cloud providers, industry associations",
#                     f"Technology Roadmap: AI enhancement, platform scalability, security certifications",
#                     f"Market Expansion: Adjacent industry assessment, international market evaluation",
#                     f"Exit Strategy: Position for strategic acquisition or IPO by 2027 with $500M+ valuation"
#                 ],
#                 "chart_type": "roi_projection",
#                 "chart_title": "Investment vs ROI Projection"
#             }
#         ]
        
#         return slides

#     def generate_ppt_data(self, slides_data, filters, competitor_landscape=None):
#         """Generate professional PowerPoint with embedded charts"""
#         try:
#             logger.info(f"Starting PPT generation with {len(slides_data)} slides")
#             logger.info(f"Filters: {filters}")
            
#             # Generate comprehensive content if not provided or insufficient
#             if not slides_data or len(slides_data) < 5:
#                 logger.warning("Insufficient slides data, generating comprehensive content")
#                 slides_data = self.generate_comprehensive_content(filters)
            
#             # Ensure chart specifications are present
#             if not any(slide.get('chart_type') for slide in slides_data):
#                 logger.warning("No chart types found in slides, adding chart specifications")
#                 self.add_chart_specifications(slides_data, filters)
            
#             prs = Presentation()
            
#             # Set slide dimensions for widescreen
#             prs.slide_width = Inches(13.33)
#             prs.slide_height = Inches(7.5)
            
#             industry_raw = filters.get('industry', 'Technology')
#             industry = self.normalize_industry_name(industry_raw)
#             logger.info(f"Industry normalized from '{industry_raw}' to '{industry}'")
            
#             for i, slide_data in enumerate(slides_data):
#                 logger.info(f"Processing slide {i+1}: {slide_data.get('title', 'No title')}")
                
#                 # Use blank layout for more control
#                 slide_layout = prs.slide_layouts[6]  # Blank layout
#                 slide = prs.slides.add_slide(slide_layout)
                
#                 # Add title
#                 title_left = Inches(0.5)
#                 title_top = Inches(0.5)
#                 title_width = Inches(12.33)
#                 title_height = Inches(1)
                
#                 title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
#                 title_frame = title_box.text_frame
#                 title_frame.text = slide_data['title']
#                 title_paragraph = title_frame.paragraphs[0]
#                 title_paragraph.font.size = Pt(28)
#                 title_paragraph.font.bold = True
#                 title_paragraph.font.color.rgb = self.colors['primary']
                
#                 # Determine layout based on chart presence
#                 has_chart = slide_data.get('chart_type') is not None
#                 logger.info(f"Slide {i+1} has_chart: {has_chart}, chart_type: {slide_data.get('chart_type')}")
                
#                 if has_chart:
#                     # Content on left, chart on right
#                     content_left = Inches(0.5)
#                     content_top = Inches(1.8)
#                     content_width = Inches(6.5)
#                     content_height = Inches(5)
                    
#                     chart_left = Inches(7.2)
#                     chart_top = Inches(1.8)
#                     chart_width = Inches(5.8)
#                     chart_height = Inches(5)
#                 else:
#                     # Full width content
#                     content_left = Inches(0.5)
#                     content_top = Inches(1.8)
#                     content_width = Inches(12.33)
#                     content_height = Inches(5)
                
#                 # Add content
#                 content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#                 content_frame = content_box.text_frame
#                 content_frame.word_wrap = True
                
#                 for j, content_item in enumerate(slide_data['content']):
#                     if j > 0:
#                         content_frame.add_paragraph()
#                     p = content_frame.paragraphs[j]
#                     p.text = f"• {content_item}"
#                     p.font.size = Pt(14)
#                     p.font.color.rgb = self.colors['text']
#                     p.space_after = Pt(8)
                
#                 # Add chart if specified
#                 if has_chart:
#                     try:
#                         chart_buffer = None
#                         chart_type = slide_data['chart_type']
                        
#                         logger.info(f"Creating {chart_type} chart for slide {i+1}")
                        
#                         if chart_type == 'market_growth':
#                             chart_buffer = self.create_market_growth_chart(industry)
#                         elif chart_type == 'competitive_landscape':
#                             chart_buffer = self.create_competitive_landscape_chart(industry)
#                         elif chart_type == 'roi_projection':
#                             chart_buffer = self.create_roi_projection_chart(industry)
#                         elif chart_type == 'market_share_pie':
#                             chart_buffer = self.create_market_share_pie_chart(industry)
                        
#                         if chart_buffer:
#                             slide.shapes.add_picture(chart_buffer, chart_left, chart_top, chart_width, chart_height)
#                             logger.info(f"✅ Successfully added {chart_type} chart to slide {i+1}")
#                         else:
#                             logger.error(f"❌ Chart buffer is None for {chart_type} on slide {i+1}")
                            
#                     except Exception as e:
#                         logger.error(f"❌ Failed to add chart to slide {i+1}: {str(e)}")
#                         import traceback
#                         logger.error(traceback.format_exc())
            
#             # Save presentation to bytes
#             ppt_buffer = io.BytesIO()
#             prs.save(ppt_buffer)
#             ppt_buffer.seek(0)
            
#             # Encode to base64
#             ppt_data = base64.b64encode(ppt_buffer.getvalue()).decode('utf-8')
            
#             # Calculate file size
#             file_size_mb = len(ppt_buffer.getvalue()) / (1024 * 1024)
            
#             charts_count = sum(1 for slide in slides_data if slide.get('chart_type'))
            
#             logger.info(f"Successfully generated PPT with {len(slides_data)} slides and {charts_count} charts")
            
#             return {
#                 "status": "success",
#                 "ppt_data": ppt_data,
#                 "filename": f"Professional_{filters.get('industry', 'Business')}_Strategy_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
#                 "file_size_mb": round(file_size_mb, 2),
#                 "slide_count": len(slides_data),
#                 "charts_included": charts_count
#             }
            
#         except Exception as e:
#             logger.error(f"Error generating PPT: {str(e)}")
#             import traceback
#             logger.error(traceback.format_exc())
#             return {
#                 "status": "error",
#                 "error_message": str(e)
#             }

#     def add_chart_specifications(self, slides_data, filters):
#         """Add chart specifications to existing slides"""
#         logger.info("Adding chart specifications to slides")
        
#         if len(slides_data) >= 5:
#             slides_data[0]["chart_type"] = "market_growth"
#             slides_data[0]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Growth"
            
#             slides_data[2]["chart_type"] = "market_share_pie"
#             slides_data[2]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Share"
            
#             slides_data[3]["chart_type"] = "competitive_landscape"
#             slides_data[3]["chart_title"] = f"{filters.get('industry', 'Technology')} Competitive Analysis"
            
#             slides_data[4]["chart_type"] = "roi_projection"
#             slides_data[4]["chart_title"] = "Investment vs ROI Projection"
            
#             logger.info("Successfully added chart specifications to slides")
#         else:
#             logger.warning(f"Not enough slides ({len(slides_data)}) to add chart specifications")
import io
import base64
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib
matplotlib.use('Agg')  # ⚠️ CRITICAL: Prevents macOS threading crashes
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import seaborn as sns
from datetime import datetime, timedelta
import random
import json
import re
import requests
from bs4 import BeautifulSoup
from market_data_collector import MarketDataCollector

# Configure matplotlib for better chart quality
plt.rcParams['figure.facecolor'] = 'white'
plt.rcParams['axes.facecolor'] = 'white'
plt.rcParams['font.size'] = 12
plt.rcParams['font.family'] = 'sans-serif'

# Set seaborn style
sns.set_style("whitegrid")
sns.set_palette("husl")

logger = logging.getLogger(__name__)

class EnhancedPPTGenerator:
    def __init__(self):
        self.colors = {
            'primary': RGBColor(37, 99, 235),      # Blue
            'secondary': RGBColor(124, 58, 237),    # Purple  
            'accent': RGBColor(5, 150, 105),       # Green
            'warning': RGBColor(217, 119, 6),      # Orange
            'danger': RGBColor(220, 38, 38),       # Red
            'dark': RGBColor(31, 41, 55),          # Dark Gray
            'light': RGBColor(249, 250, 251),      # Light Gray
            'text': RGBColor(55, 65, 81)           # Text Gray
        }
        
        # Chart color palette
        self.chart_colors = ['#2563EB', '#7C3AED', '#059669', '#D97706', '#DC2626', '#0891B2', '#C2410C']
        
        # Industry-specific comprehensive data (fallback only)
        self.industry_data = {
            'Technology': {
                'market_size': 4500,
                'growth_rate': 12.5,
                'key_players': ['Microsoft', 'Google', 'Amazon', 'Apple', 'Meta'],
                'market_shares': [22, 18, 15, 12, 8],
                'growth_rates': [8, 15, 12, 6, 25],
                'revenue_data': [100, 120, 145, 175, 210, 250, 300, 360],
                'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
                'roi_timeline': [2024, 2025, 2026, 2027],
                'roi_values': [-10, 25, 85, 180],
                'investment_values': [5, 8, 12, 15]
            },
            'Healthcare': {
                'market_size': 3200,
                'growth_rate': 8.7,
                'key_players': ['Johnson & Johnson', 'Pfizer', 'Roche', 'Novartis', 'Merck'],
                'market_shares': [15, 12, 10, 9, 8],
                'growth_rates': [6, 8, 12, 15, 10],
                'revenue_data': [80, 95, 110, 125, 140, 160, 180, 205],
                'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
                'roi_timeline': [2024, 2025, 2026, 2027],
                'roi_values': [-15, 20, 65, 150],
                'investment_values': [8, 12, 18, 25]
            },
            'Financial Services': {
                'market_size': 2800,
                'growth_rate': 6.3,
                'key_players': ['JPMorgan Chase', 'Bank of America', 'Wells Fargo', 'Goldman Sachs', 'Morgan Stanley'],
                'market_shares': [18, 15, 12, 10, 8],
                'growth_rates': [4, 6, 8, 12, 15],
                'revenue_data': [90, 100, 110, 120, 130, 145, 160, 175],
                'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
                'roi_timeline': [2024, 2025, 2026, 2027],
                'roi_values': [-8, 15, 45, 120],
                'investment_values': [6, 10, 15, 20]
            }
        }
        
        # Store current data being used
        self.current_data = None

    def normalize_industry_name(self, industry_name):
        """Normalize industry names to match data keys"""
        mapping = {
            'financial services': 'Financial Services',
            'Financial services': 'Financial Services',
            'financial': 'Financial Services',
            'technology': 'Technology',
            'tech': 'Technology',
            'healthcare': 'Healthcare',
            'health': 'Healthcare'
        }
        return mapping.get(industry_name.lower(), 'Technology')

    def extract_competitors_from_landscape(self, competitor_landscape):
        """Extract competitor data from the landscape function output"""
        companies = []
        strategies = []
        
        if not competitor_landscape or not competitor_landscape.get('landscape'):
            return companies, strategies
        
        landscape_text = competitor_landscape['landscape']
        
        # Extract company names from the landscape
        company_patterns = [
            r'\*\*([A-Z][a-zA-Z\s&]+)\*\*:',
            r'([A-Z][a-zA-Z\s&]+):',
            r'\*\*([A-Z][a-zA-Z\s&]+)\*\*'
        ]
        
        for pattern in company_patterns:
            matches = re.findall(pattern, landscape_text)
            for match in matches:
                clean_company = match.strip()
                if len(clean_company) > 2 and clean_company not in companies:
                    companies.append(clean_company)
        
        # Extract strategies
        strategy_section = re.search(r'COMPETITIVE STRATEGIES.*?(?=IBM POSITIONING|$)', 
                                    landscape_text, re.DOTALL | re.IGNORECASE)
        if strategy_section:
            strategies = re.findall(r'[-•]\s*(.+)', strategy_section.group(0))
        
        return companies[:5], [s.strip() for s in strategies if s.strip()]

    def generate_revenue_projections(self, market_size, growth_rate):
        """Generate revenue projections based on market size and growth rate"""
        base_revenue = market_size * 0.8  # Start from 80% of current market
        projections = []
        
        for i in range(8):  # 2020-2027
            revenue = base_revenue * (1 + growth_rate/100) ** i
            projections.append(round(revenue, 1))
        
        return projections

    def fetch_dynamic_market_data(self, industry, use_case, region, competitor_landscape=None):
        """Fetch real market data from web sources and competitor landscape"""
        try:
            # Extract companies from competitor landscape
            companies, strategies = self.extract_competitors_from_landscape(competitor_landscape)
            
            if not companies:
                companies = ['Microsoft', 'Google', 'Amazon', 'Apple', 'Meta']  # fallback
            
            logger.info(f"Extracted companies from landscape: {companies}")
            
            # Try to use MarketDataCollector for comprehensive data
            try:
                collector = MarketDataCollector()
                market_data = collector.fetch_market_data_from_web(industry, companies, use_case, region)
                
                if market_data and market_data.get('companies'):
                    logger.info("✅ Using MarketDataCollector for comprehensive data")
                    return market_data
                    
            except Exception as e:
                logger.warning(f"MarketDataCollector failed: {str(e)}")
            
            # Fallback to simple web scraping
            logger.info("Using fallback web scraping for market data")
            
            # Simple web scraping for market data
            search_query = f"{industry} {use_case} market size 2024"
            url = f"https://www.google.com/search?q={search_query.replace(' ', '+')}"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            market_size = None
            growth_rate = None
            
            try:
                response = requests.get(url, headers=headers, timeout=10)
                if response.status_code == 200:
                    soup = BeautifulSoup(response.content, 'html.parser')
                    text = soup.get_text()
                    
                    # Extract market size
                    market_size_matches = re.findall(r'\$(\d+(?:\.\d+)?)\s*(?:billion|B)', text, re.IGNORECASE)
                    if market_size_matches:
                        market_size = float(market_size_matches[0])
                    
                    # Extract growth rate
                    growth_matches = re.findall(r'(\d+(?:\.\d+)?)%\s*(?:CAGR|growth)', text, re.IGNORECASE)
                    if growth_matches:
                        growth_rate = float(growth_matches[0])
            
            except Exception as e:
                logger.warning(f"Web scraping failed: {str(e)}")
            
            # Create dynamic data structure
            dynamic_data = {
                'market_size': market_size or 2500,
                'growth_rate': growth_rate or 8.5,
                'key_players': companies,
                'market_shares': [20, 15, 12, 10, 8][:len(companies)],  # Estimated based on typical distribution
                'growth_rates': [8, 12, 15, 10, 6][:len(companies)],   # Estimated
                'revenue_data': self.generate_revenue_projections(market_size or 2500, growth_rate or 8.5),
                'years': [2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027],
                'roi_timeline': [2024, 2025, 2026, 2027],
                'roi_values': [-10, 25, 85, 180],
                'investment_values': [5, 8, 12, 15],
                'strategies': strategies
            }
            
            logger.info(f"✅ Created dynamic data: Market size ${market_size}B, Growth {growth_rate}%, Companies: {len(companies)}")
            return dynamic_data
            
        except Exception as e:
            logger.error(f"Error fetching dynamic market data: {str(e)}")
            return None

    def extract_data_from_content(self, slides_data, filters):
        """Extract chart data from actual slide content"""
        extracted_data = {
            'companies': [],
            'market_values': [],
            'growth_rates': [],
            'years': [],
            'revenue_data': []
        }
        
        # Parse content for numbers, companies, percentages
        for slide in slides_data:
            content_text = ' '.join(slide.get('content', []))
            
            # Extract company names (look for patterns like "Company Name leads with X%")
            company_patterns = re.findall(r'([A-Z][a-z]+ [A-Z][a-z]+|[A-Z][a-z]+) leads with (\d+)%', content_text)
            for company, percentage in company_patterns:
                if company not in extracted_data['companies']:
                    extracted_data['companies'].append(company)
                    extracted_data['market_values'].append(int(percentage))
            
            # Extract growth rates, years, revenue figures
            growth_matches = re.findall(r'(\d+(?:\.\d+)?)% growth', content_text)
            revenue_matches = re.findall(r'\$(\d+(?:\.\d+)?)B', content_text)
            year_matches = re.findall(r'(20\d{2})', content_text)
            
            extracted_data['growth_rates'].extend([float(x) for x in growth_matches])
            extracted_data['revenue_data'].extend([float(x) for x in revenue_matches])
            extracted_data['years'].extend([int(x) for x in year_matches])
        
        return extracted_data

    def create_market_growth_chart(self, industry, data=None):
        """Create a professional market growth chart with dynamic or fallback data"""
        try:
            if data is None:
                data = self.industry_data.get(industry, self.industry_data['Technology'])
            
            fig, ax = plt.subplots(figsize=(12, 8))
            fig.patch.set_facecolor('white')
            
            years = data['years']
            revenue = data['revenue_data']
            
            # Main growth line
            ax.plot(years, revenue, linewidth=4, color=self.chart_colors[0], 
                   marker='o', markersize=8, label='Market Growth')
            
            # Fill area under curve
            ax.fill_between(years, revenue, alpha=0.2, color=self.chart_colors[0])
            
            # Add trend line
            z = np.polyfit(range(len(years)), revenue, 1)
            p = np.poly1d(z)
            ax.plot(years, p(range(len(years))), "--", color=self.chart_colors[3], 
                   linewidth=3, alpha=0.8, label='Trend Line')
            
            # Highlight projection period
            ax.axvline(x=2023.5, color='gray', linestyle=':', alpha=0.7, linewidth=2)
            ax.text(2021.5, max(revenue)*0.9, 'Historical', fontsize=12, 
                   bbox=dict(boxstyle="round,pad=0.3", facecolor='lightblue', alpha=0.7))
            ax.text(2025, max(revenue)*0.9, 'Projected', fontsize=12,
                   bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgreen', alpha=0.7))
            
            # Add data labels
            for i, (year, val) in enumerate(zip(years, revenue)):
                if i % 2 == 0:  # Show every other label to avoid crowding
                    ax.annotate(f'${val}B', (year, val), textcoords="offset points", 
                              xytext=(0,10), ha='center', fontweight='bold')
            
            # Styling
            ax.set_title(f'{industry} Market Growth Trajectory (2020-2027)', 
                        fontsize=16, fontweight='bold', pad=20)
            ax.set_xlabel('Year', fontsize=14, fontweight='bold')
            ax.set_ylabel('Market Value (Billions USD)', fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3)
            ax.legend(fontsize=12)
            
            # Remove top and right spines
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            plt.tight_layout()
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            img_buffer.seek(0)
            plt.close()
            
            logger.info(f"Successfully created market growth chart for {industry}")
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating market growth chart: {str(e)}")
            plt.close()
            return None

    def create_competitive_landscape_chart(self, industry, data=None):
        """Create a competitive landscape chart with dynamic or fallback data"""
        try:
            if data is None:
                data = self.industry_data.get(industry, self.industry_data['Technology'])
            
            fig, ax = plt.subplots(figsize=(12, 8))
            fig.patch.set_facecolor('white')
            
            companies = data['key_players']
            market_shares = data['market_shares']
            growth_rates = data['growth_rates']
            
            # Ensure we have data for all companies
            if len(market_shares) < len(companies):
                market_shares.extend([5] * (len(companies) - len(market_shares)))
            if len(growth_rates) < len(companies):
                growth_rates.extend([8] * (len(companies) - len(growth_rates)))
            
            x = np.arange(len(companies))
            width = 0.35
            
            # Create bars
            bars1 = ax.bar(x - width/2, market_shares, width, label='Market Share (%)', 
                          color=self.chart_colors[0], alpha=0.8)
            bars2 = ax.bar(x + width/2, growth_rates, width, label='Growth Rate (%)', 
                          color=self.chart_colors[1], alpha=0.8)
            
            # Add value labels on bars
            for bar in bars1:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       f'{height}%', ha='center', va='bottom', fontweight='bold')
            
            for bar in bars2:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       f'{height}%', ha='center', va='bottom', fontweight='bold')
            
            # Styling
            ax.set_title(f'{industry} Competitive Analysis: Market Share vs Growth Rate', 
                        fontsize=16, fontweight='bold', pad=20)
            ax.set_xlabel('Companies', fontsize=14, fontweight='bold')
            ax.set_ylabel('Percentage', fontsize=14, fontweight='bold')
            ax.set_xticks(x)
            ax.set_xticklabels(companies, rotation=45, ha='right')
            ax.legend(fontsize=12)
            ax.grid(True, alpha=0.3)
            
            # Remove top and right spines
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            plt.tight_layout()
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            img_buffer.seek(0)
            plt.close()
            
            logger.info(f"Successfully created competitive landscape chart for {industry}")
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating competitive landscape chart: {str(e)}")
            plt.close()
            return None

    def create_roi_projection_chart(self, industry, data=None):
        """Create ROI projection chart with dynamic or fallback data"""
        try:
            if data is None:
                data = self.industry_data.get(industry, self.industry_data['Technology'])
            
            fig, ax1 = plt.subplots(figsize=(12, 8))
            fig.patch.set_facecolor('white')
            
            timeline = data['roi_timeline']
            roi_values = data['roi_values']
            investment_values = data['investment_values']
            
            # ROI line chart
            color1 = self.chart_colors[2]
            ax1.set_xlabel('Year', fontsize=14, fontweight='bold')
            ax1.set_ylabel('ROI (%)', color=color1, fontsize=14, fontweight='bold')
            line1 = ax1.plot(timeline, roi_values, 'o-', linewidth=4, color=color1, 
                           markersize=10, label='ROI %')
            ax1.tick_params(axis='y', labelcolor=color1)
            ax1.axhline(y=0, color='gray', linestyle='-', alpha=0.3)
            
            # Investment bar chart
            ax2 = ax1.twinx()
            color2 = self.chart_colors[0]
            ax2.set_ylabel('Investment ($M)', color=color2, fontsize=14, fontweight='bold')
            bars = ax2.bar(timeline, investment_values, alpha=0.6, color=color2, 
                          width=0.6, label='Investment ($M)')
            ax2.tick_params(axis='y', labelcolor=color2)
            
            # Add value labels
            for year, roi in zip(timeline, roi_values):
                ax1.annotate(f'{roi}%', (year, roi), textcoords="offset points", 
                           xytext=(0,15), ha='center', fontweight='bold', color=color1)
            
            for year, inv in zip(timeline, investment_values):
                ax2.annotate(f'${inv}M', (year, inv), textcoords="offset points", 
                           xytext=(0,5), ha='center', fontweight='bold', color=color2)
            
            # Styling
            ax1.set_title('Investment vs ROI Projection (2024-2027)', 
                         fontsize=16, fontweight='bold', pad=20)
            ax1.grid(True, alpha=0.3)
            
            plt.tight_layout()
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            img_buffer.seek(0)
            plt.close()
            
            logger.info(f"Successfully created ROI projection chart for {industry}")
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating ROI projection chart: {str(e)}")
            plt.close()
            return None

    def create_market_share_pie_chart(self, industry, data=None):
        """Create market share pie chart with dynamic or fallback data"""
        try:
            if data is None:
                data = self.industry_data.get(industry, self.industry_data['Technology'])
            
            fig, ax = plt.subplots(figsize=(10, 10))
            fig.patch.set_facecolor('white')
            
            companies = data['key_players']
            market_shares = data['market_shares']
            
            # Ensure we have market shares for all companies
            if len(market_shares) < len(companies):
                market_shares.extend([5] * (len(companies) - len(market_shares)))
            
            # Add "Others" category
            others_share = max(0, 100 - sum(market_shares))
            companies_with_others = companies + ['Others']
            shares_with_others = market_shares + [others_share]
            
            # Create pie chart
            wedges, texts, autotexts = ax.pie(shares_with_others, labels=companies_with_others, 
                                            autopct='%1.1f%%', startangle=90,
                                            colors=self.chart_colors[:len(companies_with_others)],
                                            explode=[0.05 if i == 0 else 0 for i in range(len(companies_with_others))])
            
            # Enhance text
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontweight('bold')
                autotext.set_fontsize(12)
            
            for text in texts:
                text.set_fontsize(11)
                text.set_fontweight('bold')
            
            ax.set_title(f'{industry} Market Share Distribution', 
                        fontsize=16, fontweight='bold', pad=20)
            
            plt.tight_layout()
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            img_buffer.seek(0)
            plt.close()
            
            logger.info(f"Successfully created market share pie chart for {industry}")
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating market share pie chart: {str(e)}")
            plt.close()
            return None

    def generate_comprehensive_content(self, filters):
        """Generate comprehensive slide content with chart specifications"""
        industry = filters.get('industry', 'Technology')
        use_case = filters.get('use_case', 'AI & ML')
        region = filters.get('region', 'Global')
        
        # Normalize industry name
        industry = self.normalize_industry_name(industry)
        data = self.current_data if self.current_data else self.industry_data.get(industry, self.industry_data['Technology'])
        
        slides = [
            {
                "title": f"Executive Summary: {industry} {use_case} Market Strategy",
                "content": [
                    f"Market Opportunity: ${data['market_size']}B global {industry.lower()} market growing at {data['growth_rate']}% CAGR",
                    f"Strategic Focus: {use_case} implementation targeting {region} markets with enterprise-first approach",
                    f"Competitive Landscape: Top 5 players control 75% market share, creating consolidation opportunities",
                    f"Investment Thesis: $15M investment over 24 months targeting 180% ROI by year 3",
                    f"Key Value Propositions: 40% faster implementation, 60% cost reduction, 95% customer satisfaction",
                    f"Market Entry Strategy: Phased rollout starting with pilot customers in Q2 2024",
                    f"Revenue Projections: $2M Year 1, $12M Year 2, $35M Year 3 with 85% gross margins",
                    f"Success Metrics: 50 enterprise clients, 8% market share, $50M ARR by 2027"
                ],
                "chart_type": "market_growth",
                "chart_title": f"{industry} Market Growth Trajectory"
            },
            {
                "title": f"Market Challenges & Pain Points",
                "content": [
                    f"Rapid technological obsolescence requiring constant innovation cycles",
                    f"Intense competition from both established players and disruptive startups",
                    f"Regulatory scrutiny around data privacy and antitrust concerns",
                    f"Talent acquisition and retention in highly competitive market",
                    f"Cybersecurity threats and infrastructure vulnerabilities",
                    f"Market fragmentation with 200+ vendors creating customer confusion",
                    f"Average implementation timeline of 18-24 months causing project delays",
                    f"Limited ROI visibility leading to 35% project abandonment rate"
                ],
                "chart_type": None,
                "chart_title": None
            },
            {
                "title": f"Key Market Opportunities",
                "content": [
                    f"AI and machine learning integration across all business functions",
                    f"Edge computing and 5G network expansion opportunities",
                    f"Digital transformation acceleration post-pandemic",
                    f"Emerging markets adoption of cloud technologies",
                    f"Sustainability tech and green computing initiatives",
                    f"Underserved mid-market segment representing $800M opportunity",
                    f"Geographic expansion potential in emerging markets",
                    f"Cross-selling opportunities with existing customer base"
                ],
                "chart_type": "market_share_pie",
                "chart_title": f"{industry} Market Share Distribution"
            },
            {
                "title": "Competitive Landscape Analysis",
                "content": [
                    f"Market Leaders: {data['key_players'][0]} leads with {data['market_shares'][0]}% share, followed by {data['key_players'][1]} ({data['market_shares'][1] if len(data['market_shares']) > 1 else 15}%)",
                    f"Growth Leaders: {data['key_players'][-1]} showing highest growth at {data['growth_rates'][-1] if data['growth_rates'] else 15}% vs industry average {data['growth_rate']}%",
                    f"Competitive Positioning: We target underserved mid-market with 35% better price-performance ratio",
                    f"Differentiation Strategy: AI-first approach vs competitors' legacy modernization focus",
                    f"Technology Advantage: Proprietary algorithms delivering 95% accuracy vs 78% industry standard",
                    f"Customer Acquisition: Average sales cycle 6 months vs industry standard 12-18 months",
                    f"Pricing Strategy: Premium positioning 15% above market justified by superior outcomes",
                    f"Partnership Approach: Strategic alliances with {data['key_players'][1] if len(data['key_players']) > 1 else 'key players'} for market access"
                ],
                "chart_type": "competitive_landscape",
                "chart_title": f"{industry} Competitive Analysis"
            },
            {
                "title": "Strategic Recommendations & Next Steps",
                "content": [
                    f"Immediate Actions (30 days): Secure Series A funding, finalize product roadmap, hire VP Engineering",
                    f"Short-term Goals (90 days): Complete MVP development, sign 3 pilot customers, establish advisory board",
                    f"Market Entry (6 months): Launch in {region}, build sales team of 15 FTEs, initiate marketing campaigns",
                    f"Scale Phase (12 months): Expand to 2 additional regions, achieve $5M ARR, establish channel partnerships",
                    f"Investment Recommendation: Approve $15M funding based on 180% projected ROI and market timing",
                    f"Resource Allocation: 60% product development, 25% sales & marketing, 15% operations",
                    f"Risk Mitigation: Diversify customer base, build IP portfolio, establish regulatory compliance framework",
                    f"Success Metrics: Track monthly ARR growth, customer acquisition cost, net promoter score",
                    f"Partnership Strategy: Engage tier-1 system integrators, cloud providers, industry associations",
                    f"Technology Roadmap: AI enhancement, platform scalability, security certifications",
                    f"Market Expansion: Adjacent industry assessment, international market evaluation",
                    f"Exit Strategy: Position for strategic acquisition or IPO by 2027 with $500M+ valuation"
                ],
                "chart_type": "roi_projection",
                "chart_title": "Investment vs ROI Projection"
            }
        ]
        
        return slides

    def generate_ppt_data(self, slides_data, filters, competitor_landscape=None):
        """Generate professional PowerPoint with dynamic charts from real data"""
        try:
            logger.info(f"Starting PPT generation with {len(slides_data)} slides")
            logger.info(f"Competitor landscape available: {bool(competitor_landscape)}")
            
            # Get dynamic market data
            industry_raw = filters.get('industry', 'Technology')
            industry = self.normalize_industry_name(industry_raw)
            use_case = filters.get('use_case', 'AI & ML')
            region = filters.get('region', 'Global')
            
            # Fetch real market data
            dynamic_data = self.fetch_dynamic_market_data(industry, use_case, region, competitor_landscape)
            
            if dynamic_data and dynamic_data.get('key_players'):
                logger.info("✅ Using dynamic market data for charts")
                self.current_data = dynamic_data
                data_source = "DYNAMIC"
            else:
                logger.warning("⚠️ Using fallback hardcoded data for charts")
                self.current_data = self.industry_data.get(industry, self.industry_data['Technology'])
                data_source = "HARDCODED"
            
            logger.info(f"📊 Data source: {data_source}")
            logger.info(f"📊 Market size: ${self.current_data['market_size']}B")
            logger.info(f"📊 Key players: {self.current_data['key_players']}")
            
            # Generate comprehensive content if not provided or insufficient
            if not slides_data or len(slides_data) < 5:
                logger.warning("Insufficient slides data, generating comprehensive content")
                slides_data = self.generate_comprehensive_content(filters)
            
            # Ensure chart specifications are present
            if not any(slide.get('chart_type') for slide in slides_data):
                logger.warning("No chart types found in slides, adding chart specifications")
                self.add_chart_specifications(slides_data, filters)
            
            prs = Presentation()
            
            # Set slide dimensions for widescreen
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            for i, slide_data in enumerate(slides_data):
                logger.info(f"Processing slide {i+1}: {slide_data.get('title', 'No title')}")
                
                # Use blank layout for more control
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_left = Inches(0.5)
                title_top = Inches(0.5)
                title_width = Inches(12.33)
                title_height = Inches(1)
                
                title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_frame = title_box.text_frame
                title_frame.text = slide_data['title']
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(28)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = self.colors['primary']
                
                # Determine layout based on chart presence
                has_chart = slide_data.get('chart_type') is not None
                logger.info(f"Slide {i+1} has_chart: {has_chart}, chart_type: {slide_data.get('chart_type')}")
                
                if has_chart:
                    # Content on left, chart on right
                    content_left = Inches(0.5)
                    content_top = Inches(1.8)
                    content_width = Inches(6.5)
                    content_height = Inches(5)
                    
                    chart_left = Inches(7.2)
                    chart_top = Inches(1.8)
                    chart_width = Inches(5.8)
                    chart_height = Inches(5)
                else:
                    # Full width content
                    content_left = Inches(0.5)
                    content_top = Inches(1.8)
                    content_width = Inches(12.33)
                    content_height = Inches(5)
                
                # Add content
                content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for j, content_item in enumerate(slide_data['content']):
                    if j > 0:
                        content_frame.add_paragraph()
                    p = content_frame.paragraphs[j]
                    p.text = f"• {content_item}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = self.colors['text']
                    p.space_after = Pt(8)
                
                # Add chart if specified
                if has_chart:
                    try:
                        chart_buffer = None
                        chart_type = slide_data['chart_type']
                        
                        logger.info(f"Creating {chart_type} chart for slide {i+1}")
                        
                        # Use current_data (dynamic or fallback) for all charts
                        if chart_type == 'market_growth':
                            chart_buffer = self.create_market_growth_chart(industry, self.current_data)
                        elif chart_type == 'competitive_landscape':
                            chart_buffer = self.create_competitive_landscape_chart(industry, self.current_data)
                        elif chart_type == 'roi_projection':
                            chart_buffer = self.create_roi_projection_chart(industry, self.current_data)
                        elif chart_type == 'market_share_pie':
                            chart_buffer = self.create_market_share_pie_chart(industry, self.current_data)
                        
                        if chart_buffer:
                            slide.shapes.add_picture(chart_buffer, chart_left, chart_top, chart_width, chart_height)
                            logger.info(f"✅ Successfully added {chart_type} chart to slide {i+1}")
                        else:
                            logger.error(f"❌ Chart buffer is None for {chart_type} on slide {i+1}")
                            
                    except Exception as e:
                        logger.error(f"❌ Failed to add chart to slide {i+1}: {str(e)}")
                        import traceback
                        logger.error(traceback.format_exc())
            
            # Save presentation to bytes
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            # Encode to base64
            ppt_data = base64.b64encode(ppt_buffer.getvalue()).decode('utf-8')
            
            # Calculate file size
            file_size_mb = len(ppt_buffer.getvalue()) / (1024 * 1024)
            
            charts_count = sum(1 for slide in slides_data if slide.get('chart_type'))
            
            logger.info(f"Successfully generated PPT with {len(slides_data)} slides and {charts_count} charts")
            logger.info(f"Final data source used: {data_source}")
            
            return {
                "status": "success",
                "ppt_data": ppt_data,
                "filename": f"Professional_{filters.get('industry', 'Business')}_Strategy_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                "file_size_mb": round(file_size_mb, 2),
                "slide_count": len(slides_data),
                "charts_included": charts_count,
                "data_source": data_source
            }
            
        except Exception as e:
            logger.error(f"Error generating PPT: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return {
                "status": "error",
                "error_message": str(e)
            }

    def add_chart_specifications(self, slides_data, filters):
        """Add chart specifications to existing slides"""
        logger.info("Adding chart specifications to slides")
        
        if len(slides_data) >= 5:
            slides_data[0]["chart_type"] = "market_growth"
            slides_data[0]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Growth"
            
            slides_data[2]["chart_type"] = "market_share_pie"
            slides_data[2]["chart_title"] = f"{filters.get('industry', 'Technology')} Market Share"
            
            slides_data[3]["chart_type"] = "competitive_landscape"
            slides_data[3]["chart_title"] = f"{filters.get('industry', 'Technology')} Competitive Analysis"
            
            slides_data[4]["chart_type"] = "roi_projection"
            slides_data[4]["chart_title"] = "Investment vs ROI Projection"
            
            logger.info("Successfully added chart specifications to slides")
        else:
            logger.warning(f"Not enough slides ({len(slides_data)}) to add chart specifications")
